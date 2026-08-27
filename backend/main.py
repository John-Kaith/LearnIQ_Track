import asyncio
import base64
import hashlib
import hmac
import io
import json
import os
import re
import shutil
import time
import uuid
import zipfile
from datetime import datetime, timezone
from pathlib import Path
import requests
from dotenv import load_dotenv
from fastapi import Body, FastAPI, File, UploadFile, HTTPException, Header, Query, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, RedirectResponse, Response
from fastapi.staticfiles import StaticFiles
from pypdf import PdfReader
from supabase import create_client, Client

# Always load `backend/.env` before other local modules so stale OS SMTP_USER
# / API_KEY values cannot win. `override=True` is required on Windows.
_backend_env = Path(__file__).resolve().parent / ".env"
load_dotenv(_backend_env, override=True)

import db_supabase
import email_service
import immersion_upload
from supabase_client import is_configured

if not email_service.smtp_configured():
    print(
        "[LearnIQ] Email SMTP not configured — add SMTP_USER and SMTP_PASSWORD "
        f"to {_backend_env} (Gmail App Password). Registration emails will not send."
    )
else:
    print(
        "[LearnIQ] Email SMTP ready "
        f"SMTP_USER={os.getenv('SMTP_USER')!r} SMTP_FROM={os.getenv('SMTP_FROM')!r}"
    )

BASE_DIR = Path(__file__).resolve().parent.parent
FRONTEND_DIR = BASE_DIR / "frontend"
UPLOADS_DIR = immersion_upload.UPLOAD_ROOT
LESSON_UPLOADS_DIR = UPLOADS_DIR / "lessons"
API_KEY = os.getenv("API_KEY")

AI_GENERATION_COOLDOWN_SEC = 30
_ai_generation_cooldown_until: dict[str, float] = {}

IMMERSION_CAPTURE_MAX_SKEW_MINUTES = 15


def _ai_gen_cooldown_key(kind: str, file_id: str) -> str:
    return f"{kind}:{str(file_id or '').strip()}"


def check_ai_generation_cooldown(kind: str, file_id: str) -> JSONResponse | None:
    key = _ai_gen_cooldown_key(kind, file_id)
    until = _ai_generation_cooldown_until.get(key, 0.0)
    now = time.time()
    if now < until:
        remaining = max(1, int(until - now + 0.999))
        return JSONResponse(
            {"error": f"Please wait {remaining}s before generating {kind} again."},
            status_code=429,
        )
    return None


def start_ai_generation_cooldown(kind: str, file_id: str) -> None:
    key = _ai_gen_cooldown_key(kind, file_id)
    _ai_generation_cooldown_until[key] = time.time() + AI_GENERATION_COOLDOWN_SEC

# Supabase Auth client
supabase_url = os.getenv("SUPABASE_URL")
supabase_key = os.getenv("SUPABASE_KEY")
supabase: Client = create_client(supabase_url, supabase_key)

app = FastAPI(title="LearnIQ Track API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def require_supabase():
    if not is_configured():
        return JSONResponse(
            {
                "error": "Supabase is not configured. Add SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY "
                "(or SUPABASE_KEY) to backend/.env and run supabase_schema.sql in the Supabase SQL editor."
            },
            status_code=503,
        )
    return None


def student_id_number_from_authorization(authorization: str | None) -> str | None:
    """Resolve school id_number from Supabase access token (Authorization: Bearer …)."""
    if not authorization or not authorization.lower().startswith("bearer "):
        return None
    token = authorization.split(" ", 1)[1].strip()
    if not token:
        return None
    try:
        ures = supabase.auth.get_user(jwt=token)
        if not ures:
            return None
        user = getattr(ures, "user", None)
        email = getattr(user, "email", None) if user is not None else None
        if not email:
            return None
        prof = db_supabase.get_profile_by_email(email)
        if not prof:
            return None
        sid = db_supabase.profile_lrn_from_row(prof)
        return str(sid).strip() if sid else None
    except Exception as e:
        print("student_id_number_from_authorization:", e)
        return None


def resolve_student_id_number_or_403(body: dict, authorization: str | None):
    """Returns (id_number, None) or (None, JSONResponse). A valid session
    token is always required now — a request used to be able to skip
    signing in entirely and just declare a student_id_number in the body,
    which meant anyone could act as any student without ever logging in.
    (Confirmed every current caller already sends a token; see
    backend_auth_audit_phases memory / Phase 3 for the audit.)"""
    token_sid = student_id_number_from_authorization(authorization)
    if not token_sid:
        return None, JSONResponse({"error": "Sign in required."}, status_code=401)
    body_sid = (body.get("student_id_number") or body.get("student_id") or "").strip()
    if body_sid and body_sid != token_sid:
        return None, JSONResponse(
            {"error": "student_id_number does not match signed-in user."},
            status_code=403,
        )
    return token_sid, None


def _can_view_student_data(
    authorization: str | None, target_student_id_number: str
) -> tuple[bool, str | None, JSONResponse | None]:
    """True if the caller may view/act on this student's data: the student
    themselves, an admin, or a teacher who has this student enrolled in one
    of their subjects. Returns (allowed, caller_idn, error_response)."""
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return False, None, JSONResponse({"error": "Sign in required."}, status_code=401)
    target = str(target_student_id_number or "").strip()
    if caller_idn == target:
        return True, caller_idn, None
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    if role == "admin":
        return True, caller_idn, None
    if role == "teacher":
        target_prof = db_supabase.get_profile_by_id_number(target)
        target_uuid = str((target_prof or {}).get("id") or "")
        if target_uuid and target_uuid in db_supabase._student_uuids_enrolled_in_teacher_subjects(caller_idn):
            return True, caller_idn, None
    return False, caller_idn, JSONResponse(
        {"error": "You don't have access to this student's data."}, status_code=403
    )


def _require_signed_in(authorization: str | None) -> JSONResponse | None:
    """Any signed-in student/teacher/admin — just blocks anonymous callers
    from spending the AI generation budget. Returns an error response, or
    None if the caller is signed in."""
    if not student_id_number_from_authorization(authorization):
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    return None


def _can_view_teacher_data(authorization: str | None, target_teacher_id_number: str) -> tuple[bool, JSONResponse | None]:
    """True if the caller may view this teacher's own data: that teacher, or an admin."""
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return False, JSONResponse({"error": "Sign in required."}, status_code=401)
    target = str(target_teacher_id_number or "").strip()
    if caller_idn == target:
        return True, None
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    if role == "admin":
        return True, None
    return False, JSONResponse({"error": "You don't have access to this teacher's data."}, status_code=403)


def gemini_text_from_result(result: dict) -> str:
    try:
        parts = result["candidates"][0]["content"]["parts"]
        if not parts:
            return ""
        return parts[0].get("text") or ""
    except (KeyError, IndexError, TypeError):
        return ""


def strip_outer_markdown_code_fence(text: str) -> str:
    """Remove accidental ```markdown … ``` wrapper from model output."""
    t = (text or "").strip()
    if not t.startswith("```"):
        return t
    lines = t.split("\n")
    if lines and lines[0].strip().startswith("```"):
        lines = lines[1:]
    while lines and lines[-1].strip() == "```":
        lines = lines[:-1]
    return "\n".join(lines).strip()


REVIEWER_SOURCE_MAX_CHARS = 120_000
LESSON_EXTRACT_MAX_CHARS = 120_000
LESSON_VISION_MAX_SLIDES = 15
# Whole-file base64 is stored in Postgres; keep uploads reasonable for insert + JSON latency.
LESSON_UPLOAD_MAX_BYTES = 500 * 1024 * 1024
LESSON_TEXT_EXTRACT_TIMEOUT_SEC = 120

REVIEWER_PROMPT_TEMPLATE = (
    "You are an experienced teacher writing a study reviewer for senior high school students.\n\n"
    "TASK: Using ONLY the lesson text below, produce a clear reviewer students can use to study.\n\n"
    "OUTPUT FORMAT (GitHub-flavored Markdown only; no HTML tags):\n"
    "1) Start with exactly one level-1 heading on its own line: # <short topic title>\n"
    "2) One short opening paragraph (3–5 sentences) that explains the big idea in plain language.\n"
    "3) ## Key Points — then 4–7 bullets. Each bullet must be its own line starting with \"• \" "
    "(bullet character + space).\n"
    "4) ## Important Definitions — one term per line using this pattern: Term — short definition "
    "(use an em dash between term and definition).\n"
    "5) ## Summary — one short closing paragraph (2–4 sentences) that ties ideas together.\n\n"
    "STYLE RULES:\n"
    "- Balance short paragraphs with lists; avoid one giant paragraph or one endless bullet list.\n"
    "- Sound clear and educational; avoid hype, filler, or stock AI phrases "
    '(e.g. avoid "delve", "leverage", "In conclusion", "It is important to note").\n'
    "- Simplify complex ideas; do not repeat the same point in multiple sections.\n"
    "- Keep vocabulary accessible; define jargon briefly when it appears.\n"
    "- Write in the same language as the lesson when the lesson is not English; otherwise use English.\n"
    "- Do not wrap the entire answer in markdown code fences (no ```).\n"
    "- Do not add links, images, URLs, or quiz questions.\n\n"
    "LESSON TEXT:\n{source}\n"
)


def parse_model_json(raw: str):
    """Parse JSON from model output; strips optional markdown code fences."""
    text = raw.strip()
    if text.startswith("```"):
        lines = text.split("\n")
        if len(lines) >= 2:
            text = "\n".join(lines[1:])
        if "```" in text:
            text = text[: text.index("```")].strip()
    return json.loads(text)


def friendly_ai_error(result: dict | str | None) -> str:
    """
    Convert provider errors into a user-friendly message (no raw payload in UI).
    """
    try:
        if isinstance(result, dict):
            msg = (
                result.get("error", {}).get("message")
                if isinstance(result.get("error"), dict)
                else result.get("message")
            )
            msg = str(msg or "").lower()
            if "high demand" in msg or "overloaded" in msg or "unavailable" in msg:
                return "AI servers are currently busy. Please try again."
            if "quota" in msg or "rate limit" in msg or "resource" in msg:
                return "AI quota is currently exceeded. Please try again later."
            if "api key" in msg or "permission" in msg or "denied" in msg:
                return "AI access is currently unavailable. Please contact the administrator."
        return "Failed to generate content. Please retry."
    except Exception:
        return "Failed to generate content. Please retry."


def normalize_quiz_questions(obj, desired_count: int):
    """
    Accept either {"questions":[...]} or a raw list, then validate/normalize.
    Enforces: question(str), choices(list of 4 strings), answer(letter A-D).
    """
    if isinstance(obj, dict) and isinstance(obj.get("questions"), list):
        questions = obj["questions"]
    elif isinstance(obj, list):
        questions = obj
    else:
        raise ValueError("Invalid quiz JSON format")

    out = []
    for q in questions:
        if not isinstance(q, dict):
            continue
        question = str(q.get("question") or "").strip()
        choices = q.get("choices") or []
        answer = q.get("answer")
        if not question:
            continue
        if not isinstance(choices, list):
            continue
        choices = [str(c) for c in choices if str(c).strip()]
        if len(choices) != 4:
            continue

        # Normalize answer to a letter (A-D).
        ans = str(answer or "").strip()
        if len(ans) == 1 and ans.upper() in ("A", "B", "C", "D"):
            ans_letter = ans.upper()
        else:
            # If model returns full choice text, try matching by prefix or exact.
            ans_letter = ""
            for idx, c in enumerate(choices):
                if ans.strip().lower() == c.strip().lower():
                    ans_letter = ("A", "B", "C", "D")[idx]
                    break
            if not ans_letter:
                first = ans[:1].upper()
                if first in ("A", "B", "C", "D"):
                    ans_letter = first
        if not ans_letter:
            continue

        out.append({"question": question, "choices": choices, "answer": ans_letter})
        if desired_count and len(out) >= desired_count:
            break

    if not out:
        raise ValueError("No valid quiz questions parsed")
    return out


def require_gemini_key():
    if not API_KEY or not str(API_KEY).strip():
        return JSONResponse(
            {"error": "Missing API_KEY in backend .env (Google Gemini). AI generation is disabled until you add it."},
            status_code=503,
        )
    return None


@app.get("/health")
def health():
    return {
        "ok": True,
        "has_api_key": bool(API_KEY and API_KEY.strip()),
        "has_supabase": is_configured(),
        "email_smtp_configured": email_service.smtp_configured(),
    }


@app.get("/test-db")
def test_db():
    """
    Quick Supabase smoke test: insert -> select -> update (then cleanup).
    Useful for confirming DB connectivity and table write/read behavior.
    """
    err = require_supabase()
    if err is not None:
        return err
    test_id = f"SBTEST-{uuid.uuid4().hex[:8]}"
    test_email = f"{test_id.lower()}@example.com"
    result = {
        "configured": is_configured(),
        "test_id_number": test_id,
        "insert_ok": False,
        "select_ok": False,
        "update_ok": False,
        "cleanup_ok": False,
    }
    try:
        inserted = db_supabase.insert_profile(
            id_number=test_id,
            email=test_email,
            password="TempPass123",
            role="student",
            last_name="Test",
            first_name="Supabase",
        )
        result["insert_ok"] = bool(
            inserted and db_supabase.profile_lrn_from_row(inserted) == test_id
        )

        rows = [
            r
            for r in db_supabase.list_profiles()
            if db_supabase.profile_lrn_from_row(r) == test_id
        ]
        result["select_ok"] = bool(rows)

        rows_after = [
            r
            for r in db_supabase.list_profiles()
            if db_supabase.profile_lrn_from_row(r) == test_id
        ]
        result["update_ok"] = bool(rows_after)

        # Keep test data out of production tables.
        db_supabase._sb().table("profiles").delete().eq(db_supabase.PROFILE_LRN_COLUMN, test_id).execute()
        rows_cleanup = [
            r
            for r in db_supabase.list_profiles()
            if db_supabase.profile_lrn_from_row(r) == test_id
        ]
        result["cleanup_ok"] = not rows_cleanup

        result["message"] = "Supabase test completed."
        return result
    except Exception as e:
        result["error"] = str(e)
        return JSONResponse(result, status_code=500)


@app.get("/")
def home():
    return RedirectResponse(url="/login.html", status_code=302)


# --- Profiles (replaces local-only auth storage for the API side) ---


@app.post("/login")
async def login_user(body: dict):
    print("LOGIN ENDPOINT HIT")
    print(f"DEBUG: Login attempt - body: {body}")
    
    err = require_supabase()
    if err is not None:
        print(f"DEBUG: Supabase connection error: {err}")
        return err
    
    try:
        password = body.get("password") or ""
        identifier = (
            body.get("identifier")
            or body.get("email")
            or body.get("lrn")
            or body.get("id_number")
            or ""
        ).strip()
        login_method = (body.get("login_method") or "").strip().lower()

        if not identifier or not password:
            return JSONResponse(
                {"error": "LRN or email and password are required."},
                status_code=400,
            )

        use_email = login_method == "email" or (
            login_method != "lrn" and "@" in identifier
        )

        if use_email:
            email = identifier.lower()
            user_profile = db_supabase.get_profile_by_email(email)
            if not user_profile:
                return JSONResponse({"error": "Invalid credentials."}, status_code=401)
        else:
            user_profile = db_supabase.get_profile_by_id_number(identifier)
            if not user_profile:
                return JSONResponse({"error": "Invalid credentials."}, status_code=401)
            email = (user_profile.get("email") or "").strip().lower()
            if not email:
                return JSONResponse(
                    {
                        "error": "No email is linked to this LRN. Contact your administrator."
                    },
                    status_code=404,
                )

        print(f"LOGIN IDENTIFIER: {identifier!r}, AUTH EMAIL: {email}")

        # Authenticate with Supabase (email + password)
        print(f"DEBUG: Attempting Supabase auth for email: {email}")
        try:
            auth_response = supabase.auth.sign_in_with_password({
                "email": email,
                "password": password
            })
            print(f"DEBUG: Supabase auth response: {type(auth_response)}")
            print(f"DEBUG: Auth user: {auth_response.user}")
            print(f"DEBUG: Auth session: {auth_response.session}")
            print(f"LOGIN PASSWORD VALIDATION RESULT: {bool(auth_response.user)}")
        except Exception as auth_error:
            print(f"DEBUG: Supabase auth exception: {auth_error}")
            print(f"DEBUG: Auth exception type: {type(auth_error)}")
            return JSONResponse({"error": "Invalid credentials."}, status_code=401)
        
        if not auth_response.user:
            print(f"DEBUG: No user returned from auth")
            return JSONResponse({"error": "Invalid credentials."}, status_code=401)
        
        # Return safe user data with auth session
        print(f"DEBUG: Preparing successful response")
        print(f"DEBUG: Raw user_profile from DB: {user_profile}")
        print(f"DEBUG: user_profile keys: {list(user_profile.keys()) if user_profile else 'None'}")
        print(f"DEBUG: user_profile.get('role'): {user_profile.get('role') if user_profile else 'None'}")
        print(f"DEBUG: Type of role: {type(user_profile.get('role')) if user_profile else 'None'}")
        
        try:
            role_value = user_profile.get("role")
            print(f"DEBUG: Role value before processing: '{role_value}'")
            print(f"DEBUG: Role value trimmed: '{role_value.strip() if role_value else None}'")
            print(f"DEBUG: Role value lowercased: '{role_value.strip().lower() if role_value else None}'")
            
            safe_user = db_supabase.serialize_public_profile(user_profile)
            safe_user["role"] = role_value.strip().lower() if role_value else "student"
            safe_user["access_token"] = auth_response.session.access_token
            safe_user["refresh_token"] = auth_response.session.refresh_token
            print(f"DEBUG: Safe user data prepared: {safe_user}")
            print(f"DEBUG: Final role in safe_user: '{safe_user['role']}'")
            print(f"DEBUG: Final role type: {type(safe_user['role'])}")
        except Exception as response_error:
            print(f"DEBUG: Response preparation exception: {response_error}")
            print(f"DEBUG: Response exception type: {type(response_error)}")
            return JSONResponse({"error": "Error preparing user response."}, status_code=500)
        
        print(f"DEBUG: Login successful for user: {email}")
        print(f"FINAL LOGIN RESPONSE: {safe_user}")
        print(f"FINAL RESPONSE STRUCTURE: {{'user': {safe_user}, 'message': 'Login successful'}}")
        print(f"LOGIN RESPONSE PAYLOAD: {{'user': {safe_user}, 'message': 'Login successful'}}")
        return {"user": safe_user, "message": "Login successful"}
        
    except Exception as e:
        print(f"DEBUG: Login endpoint exception: {e}")
        print(f"DEBUG: Exception type: {type(e)}")
        print(f"DEBUG: Exception args: {e.args}")
        import traceback
        print(f"DEBUG: Full traceback: {traceback.format_exc()}")
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/auth/refresh")
async def refresh_access_token(body: dict):
    """Exchange a Supabase refresh_token for a new access_token.

    Body: { "refresh_token": "..." }
    Returns: { "access_token": "...", "refresh_token": "...", "expires_in": int }
    """
    err = require_supabase()
    if err is not None:
        return err
    try:
        refresh_token = (body.get("refresh_token") or "").strip()
        if not refresh_token:
            return JSONResponse({"error": "refresh_token required."}, status_code=400)
        result = supabase.auth.refresh_session(refresh_token)
        session = getattr(result, "session", None)
        if not session or not getattr(session, "access_token", None):
            return JSONResponse({"error": "Could not refresh session."}, status_code=401)
        return {
            "access_token": session.access_token,
            "refresh_token": getattr(session, "refresh_token", refresh_token),
            "expires_in": getattr(session, "expires_in", None),
        }
    except Exception as e:
        msg = str(e) or "Refresh failed."
        print("AUTH REFRESH ERROR:", msg)
        return JSONResponse({"error": msg}, status_code=401)


def _reset_password_redirect_url() -> str:
    """Same host/port as APP_LOGIN_URL (whatever that's set to per-environment), pointed at reset-password.html."""
    login_url = (os.getenv("APP_LOGIN_URL") or "http://127.0.0.1:8000/login.html").strip()
    base = login_url.rsplit("/", 1)[0]
    return f"{base}/reset-password.html"


@app.post("/forgot-password")
async def forgot_password(body: dict):
    err = require_supabase()
    if err is not None:
        return err
    try:
        email = (body.get("email") or "").strip()

        if not email:
            return JSONResponse({"error": "Email is required."}, status_code=400)

        # Use Supabase Auth to send password reset email
        redirect_to = _reset_password_redirect_url()
        print(f"[DEBUG] /forgot-password redirectTo = {redirect_to!r}")
        supabase.auth.reset_password_for_email(email, {
            "redirectTo": redirect_to
        })

        return {"message": "Password reset instructions have been sent to your email."}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/reset-password")
async def reset_password(body: dict):
    """Completes a Supabase password recovery: exchanges the recovery
    token/code from the emailed link for a session, then sets the new
    password. Frontend (reset-password.html) reads the token/code from
    the URL and calls this."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        new_password = (body.get("new_password") or "").strip()
        if len(new_password) < 8:
            return JSONResponse({"error": "Password must be at least 8 characters."}, status_code=400)

        access_token = body.get("access_token")
        refresh_token = body.get("refresh_token")
        code = body.get("code")

        if access_token and refresh_token:
            supabase.auth.set_session(access_token, refresh_token)
        elif code:
            supabase.auth.exchange_code_for_session({"auth_code": code})
        else:
            return JSONResponse(
                {"error": "This reset link is invalid or missing. Request a new one."},
                status_code=400,
            )

        supabase.auth.update_user({"password": new_password})
        return {"message": "Password updated. You can now log in with your new password."}
    except Exception as e:
        return JSONResponse({"error": "Could not reset password. The link may have expired — request a new one."}, status_code=400)


@app.post("/validate-session")
async def validate_session(body: dict):
    err = require_supabase()
    if err is not None:
        return err
    try:
        access_token = body.get("access_token")
        if not access_token:
            return JSONResponse({"error": "Access token required."}, status_code=400)
        
        # Set the session and get current user
        supabase.auth.set_session(access_token)
        user = supabase.auth.get_user()
        
        if not user.user:
            return JSONResponse({"error": "Invalid session."}, status_code=401)
        
        # Get user profile from our profiles table
        user_profile = db_supabase.get_profile_by_email(user.user.email)
        if not user_profile:
            return JSONResponse({"error": "User profile not found."}, status_code=404)
        
        # Return safe user data
        safe_user = db_supabase.serialize_public_profile(user_profile)
        return {"user": safe_user, "message": "Session valid"}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


_STUDENT_STRANDS = frozenset({"ABM", "HUMSS", "STEM", "TVL-HE"})
_STUDENT_GRADE_LEVELS = frozenset({"11", "12"})


def _friendly_registration_error(exc: Exception) -> tuple[str, int]:
    """Map Supabase/Postgres errors to short messages for teachers/admins."""
    msg = str(exc)
    low = msg.lower()
    if "(lrn)=" in low or "profiles_id_number_key" in low or "profiles_lrn" in low:
        m = re.search(r"\(lrn\)=\(([^)]+)\)", msg)
        if m:
            return (
                f"LRN {m.group(1)} is already registered. Use a different LRN or ask admin to remove the old account.",
                409,
            )
        return (
            "This LRN is already registered. Use a different learner reference number.",
            409,
        )
    if "duplicate" in low and "email" in low:
        return "This email is already registered. Use a different email address.", 409
    if "duplicate" in low or "unique" in low or "23505" in msg:
        return "An account with this LRN or email already exists.", 409
    return msg, 400


def _normalize_name_suffix(raw: str) -> str | None:
    s = (raw or "").strip()
    if not s:
        return None
    if len(s) > 20:
        return None
    key = s.upper().replace(".", "")
    aliases = {
        "JR": "Jr",
        "JUNIOR": "Jr",
        "SR": "Sr",
        "SENIOR": "Sr",
        "III": "III",
        "IV": "IV",
        "V": "V",
    }
    return aliases.get(key, s)


@app.post("/register")
async def register_user(body: dict, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    try:
        id_number = (body.get("id_number") or "").strip()
        email = (body.get("email") or "").strip().lower()
        password_raw = (body.get("password") or "").strip()
        auto_generate = bool(body.get("auto_generate_password")) or not password_raw
        password = password_raw or (
            email_service.generate_registration_password()
            if auto_generate
            else ""
        )
        role = (body.get("role") or "student").strip().lower()

        if not id_number or not email:
            return JSONResponse(
                {"error": "id_number and email are required."},
                status_code=400,
            )
        if not password:
            return JSONResponse(
                {"error": "password is required, or leave it blank to auto-generate."},
                status_code=400,
            )

        if role not in ("student", "teacher"):
            role = "student"

        # Teacher accounts can only be created by an admin — no public
        # teacher self-signup. The caller must present a valid admin
        # Bearer token (admin-teacher-registration.html sends its own
        # session token); a client-supplied flag would be spoofable, so
        # this is enforced server-side via the token itself.
        if role == "teacher":
            creator_idn = student_id_number_from_authorization(authorization)
            creator_prof = db_supabase.get_profile_by_id_number(creator_idn) if creator_idn else None
            if not creator_prof or str(creator_prof.get("role") or "").strip().lower() != "admin":
                return JSONResponse(
                    {"error": "Teacher accounts can only be created by an admin."},
                    status_code=403,
                )

        last_name = (body.get("last_name") or "").strip()
        first_name = (body.get("first_name") or "").strip()
        middle_name = (body.get("middle_name") or "").strip()
        name_suffix = _normalize_name_suffix(body.get("name_suffix") or "")
        if (body.get("name_suffix") or "").strip() and not name_suffix:
            return JSONResponse(
                {"error": "Suffix is too long (max 20 characters)."},
                status_code=400,
            )
        grade_level = (body.get("grade_level") or "").strip()
        strand_raw = (body.get("strand") or "").strip().upper().replace(" ", "-")
        strand_aliases = {"HUMMS": "HUMSS", "TVLHE": "TVL-HE"}
        strand = strand_aliases.get(strand_raw, strand_raw)
        # Optional — only the admin registration form sends this today (a
        # fixed dropdown, not free text). Self-signup omits it; an admin
        # assigns it later via Manage Sections in that case.
        section = (body.get("section") or "").strip()

        if not last_name or not first_name:
            return JSONResponse(
                {"error": "last_name and first_name are required."},
                status_code=400,
            )

        if role == "student":
            if grade_level not in _STUDENT_GRADE_LEVELS:
                return JSONResponse(
                    {"error": "grade_level must be 11 or 12."},
                    status_code=400,
                )
            if strand not in _STUDENT_STRANDS:
                return JSONResponse(
                    {"error": "strand must be one of: ABM, HUMSS, STEM, TVL-HE."},
                    status_code=400,
                )
            if section:
                valid_sections = {
                    s["name"] for s in db_supabase.list_sections(grade_level, strand)
                }
                if section not in valid_sections:
                    return JSONResponse(
                        {"error": "That section is not available for this grade level and strand."},
                        status_code=400,
                    )

        auth_meta = {
            "id_number": id_number,
            "role": role,
            "last_name": last_name,
            "first_name": first_name,
            "middle_name": middle_name,
            "name_suffix": name_suffix or "",
        }
        if role == "student":
            auth_meta["grade_level"] = grade_level
            auth_meta["strand"] = strand

        auth_response = supabase.auth.sign_up({
            "email": email,
            "password": password,
            "options": {"data": auth_meta},
        })

        if not auth_response.user:
            return JSONResponse({"error": "Failed to create account."}, status_code=400)

        profile = db_supabase.insert_profile(
            id_number=id_number,
            email=email,
            password="",
            role=role,
            auth_user_id=auth_response.user.id,
            last_name=last_name,
            first_name=first_name,
            middle_name=middle_name or None,
            name_suffix=name_suffix,
            grade_level=grade_level if role == "student" else None,
            strand=strand if role == "student" else None,
            section=section if role == "student" else None,
        )
        
        profile = dict(profile)
        profile.pop("password", None)

        credentials_emailed = False
        credentials_email_error = None
        if auto_generate and not password_raw:
            sent_ok, err = await asyncio.to_thread(
                email_service.send_registration_credentials_email,
                to_email=email,
                first_name=first_name,
                last_name=last_name,
                middle_name=middle_name,
                name_suffix=name_suffix or "",
                lrn=id_number,
                login_email=email,
                password=password,
                role=role,
            )
            credentials_emailed = sent_ok
            credentials_email_error = err

        if auto_generate and not password_raw and credentials_emailed:
            message = (
                "Account created. Login credentials were emailed to "
                f"{email} (LRN/ID, email, and password)."
            )
        elif auto_generate and not password_raw:
            message = (
                "Account created with an auto-generated password, but the "
                "credentials email could not be sent."
            )
            if credentials_email_error:
                message += f" {credentials_email_error}"
        else:
            message = (
                "Account created successfully. "
                "You can sign in after confirming your email."
            )

        return {
            "user": profile,
            "message": message,
            "email": email,
            "auto_generated_password": auto_generate and not password_raw,
            "credentials_emailed": credentials_emailed,
            "credentials_email_error": credentials_email_error,
        }
    except Exception as e:
        friendly, status = _friendly_registration_error(e)
        return JSONResponse({"error": friendly}, status_code=status)


@app.get("/users")
def get_users(authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        rows = db_supabase.get_all_profiles()
        out = []
        for r in rows:
            x = dict(r)
            x.pop("password", None)
            out.append(x)
        return out
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/admin/stats")
def get_admin_dashboard_stats(authorization: str | None = Header(default=None)):
    """Aggregated metrics for the admin dashboard (Supabase)."""
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        return db_supabase.get_admin_dashboard_stats()
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/admin/attendance-logs")
def admin_list_attendance_logs(
    limit: int = Query(default=200, ge=1, le=500),
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        return {"logs": db_supabase.list_all_attendance_logs(limit)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/admin/journals-feed")
def admin_list_journals_feed(
    limit: int = Query(default=200, ge=1, le=500),
    authorization: str | None = Header(default=None),
):
    """All journal submissions (admin)."""
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        return {"journals": db_supabase.list_all_journals_admin(limit)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/admin/profile/{id_number}")
def admin_get_profile_by_id_number(id_number: str, authorization: str | None = Header(default=None)):
    """Single profile row from Supabase (admin UI / modals). Password omitted."""
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        prof = db_supabase.get_profile_by_id_number(id_number)
        if not prof:
            return JSONResponse({"error": "Profile not found."}, status_code=404)
        out = dict(prof)
        out.pop("password", None)
        return out
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/admin/recent-activity")
def admin_recent_activity(
    limit: int = Query(default=12, ge=1, le=50),
    authorization: str | None = Header(default=None),
):
    """Recent registrations, uploads, and quiz attempts for the admin dashboard."""
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        return {"items": db_supabase.get_admin_recent_activity(limit)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


# --- Lessons & AI (Supabase) ---


@app.get("/teacher/lessons")
def list_teacher_lessons(
    teacher_id_number: str = Query(...),
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    allowed, bad = _can_view_teacher_data(authorization, teacher_id_number)
    if not allowed:
        return bad
    try:
        rows = db_supabase.list_teacher_lessons(teacher_id_number)
        return {"lessons": [db_supabase.lesson_to_api_list_item(r) for r in rows]}
    except Exception as e:
        import traceback
        print("TEACHER LESSONS ERROR:", repr(e))
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/dashboard-stats")
def get_teacher_dashboard_stats(authorization: str | None = Header(default=None)):
    """Teacher LearnIQ dashboard: upload counts, students with activity, avg quiz score (Bearer token)."""
    err = require_supabase()
    if err is not None:
        return err
    idn = student_id_number_from_authorization(authorization)
    if not idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        prof = db_supabase.get_profile_by_id_number(idn)
        role = str((prof or {}).get("role") or "").strip().lower()
        if role != "teacher":
            return JSONResponse({"error": "Teacher access only."}, status_code=403)
        return db_supabase.get_teacher_learniq_dashboard_stats(idn)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)

def _profile_response_payload(prof: dict) -> dict:
    """Shared shape returned by GET /me and PATCH /me/profile."""
    return db_supabase.serialize_public_profile(prof)


@app.get("/me")
def get_me(authorization: str | None = Header(default=None)):
    """Return signed-in user's profile row (Bearer token)."""
    err = require_supabase()
    if err is not None:
        return err
    sid = student_id_number_from_authorization(authorization)
    if not sid:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        prof = db_supabase.get_profile_by_id_number(sid)
        if not prof:
            return JSONResponse({"error": "Profile not found."}, status_code=404)
        return _profile_response_payload(prof)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.patch("/me/profile")
async def patch_my_profile(
    body: dict = Body(default={}),
    authorization: str | None = Header(default=None),
):
    """Signed-in user updates their own editable profile fields.

    Accepts any subset of: bio, phone, section, dob (YYYY-MM-DD), address,
    avatar_data (data URL string). Other keys are ignored. Empty strings
    clear the column (stored as NULL).
    """
    err = require_supabase()
    if err is not None:
        return err
    sid = student_id_number_from_authorization(authorization)
    if not sid:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        payload = body if isinstance(body, dict) else {}
        editable: dict = {}
        for f in db_supabase.PROFILE_EXTRA_FIELDS:
            if f in payload:
                editable[f] = payload[f]
        if not editable:
            return JSONResponse({"error": "No editable fields provided."}, status_code=400)
        # Basic guard: keep avatar payload reasonable (~2 MB of base64 ≈ 1.5 MB image).
        av = editable.get("avatar_data")
        if isinstance(av, str) and len(av) > 2_500_000:
            return JSONResponse(
                {"error": "Avatar is too large after encoding. Try a smaller image."},
                status_code=413,
            )
        updated = db_supabase.update_profile_extras(sid, editable)
        if not updated:
            return JSONResponse({"error": "Could not update profile."}, status_code=500)
        return _profile_response_payload(updated)
    except Exception as e:
        msg = str(e)
        if "column" in msg.lower() and ("does not exist" in msg.lower() or "schema cache" in msg.lower()):
            return JSONResponse(
                {
                    "error": (
                        "Profile columns are missing in the database. Run "
                        "`backend/migrations/profile_extras.sql` in the Supabase "
                        "SQL editor and try again."
                    )
                },
                status_code=500,
            )
        return JSONResponse({"error": msg}, status_code=502)


@app.get("/lessons")
def list_all_lessons_admin():
    """All uploaded lesson files (admin Uploaded Files section + dashboard stat)."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        rows = db_supabase.list_all_lessons()
        lessons_out = []
        for r in rows:
            base = db_supabase.lesson_to_api_list_item(r)
            lessons_out.append(
                {
                    **base,
                    "file_type": r.get("file_type") or "",
                    "created_at": r.get("created_at"),
                    "teacher_id_number": r.get("teacher_id_number") or "",
                    "is_published": bool(r.get("is_published")),
                }
            )
        return {"lessons": lessons_out, "count": len(lessons_out)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/publish-lesson")
async def publish_lesson(body: dict, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    lesson_id = body.get("lesson_id") or body.get("file_id")
    if not lesson_id:
        return JSONResponse({"error": "lesson_id (or file_id) is required."}, status_code=400)
    try:
        lesson = db_supabase.get_lesson_row(str(lesson_id))
        if not lesson:
            return JSONResponse({"error": "Lesson not found."}, status_code=404)
        allowed, bad = _can_view_teacher_data(authorization, str(lesson.get("teacher_id_number") or ""))
        if not allowed:
            return bad
        db_supabase.publish_lesson(str(lesson_id))
        return {"published_file_id": lesson_id, "message": "Students can now open this lesson on their dashboard."}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/unpublish-lesson")
async def unpublish_lesson(body: dict, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    lesson_id = body.get("lesson_id") or body.get("file_id")
    if not lesson_id:
        return JSONResponse({"error": "lesson_id (or file_id) is required."}, status_code=400)
    try:
        lesson = db_supabase.get_lesson_row(str(lesson_id))
        if not lesson:
            return JSONResponse({"error": "Lesson not found."}, status_code=404)
        allowed, bad = _can_view_teacher_data(authorization, str(lesson.get("teacher_id_number") or ""))
        if not allowed:
            return bad
        # Bug fix: this used to call unpublish_all_lessons(), silently
        # unpublishing every lesson school-wide instead of just this one.
        db_supabase.unpublish_lesson(str(lesson_id))
        return {"unpublished_file_id": lesson_id, "message": "Lesson is no longer visible to students."}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/lessons")
def get_student_lessons(
    subject_id: str | None = None,
    student_id_number: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse(
            {"error": "student_id_number is required to list lessons."},
            status_code=400,
        )
    allowed, _, bad = _can_view_student_data(authorization, sid)
    if not allowed:
        return bad
    student_uuid = db_supabase.profile_uuid_for_id_number(sid)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)
    try:
        lessons = db_supabase.list_published_lessons_for_student(
            student_uuid,
            subject_id=subject_id,
        )
        print(
            "STUDENT LESSONS DEBUG: Found",
            len(lessons),
            "enrolled lessons (subject_id=",
            subject_id,
            ")",
        )
        return {"lessons": lessons}
    except Exception as e:
        print("STUDENT LESSONS ERROR:", str(e))
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/subjects")
def list_student_enrolled_subjects_endpoint(
    student_id_number: str = Query(...),
    period_id: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
):
    """Subjects the student is enrolled in (current grading period by default)."""
    err = require_supabase()
    if err is not None:
        return err
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    allowed, _, bad = _can_view_student_data(authorization, sid)
    if not allowed:
        return bad
    student_uuid = db_supabase.profile_uuid_for_id_number(sid)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)
    try:
        subjects = db_supabase.list_enrolled_subjects_for_student(student_uuid, period_id)
        return {"subjects": subjects, "count": len(subjects)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/subjects/archived")
def list_student_archived_subjects_endpoint(
    student_id_number: str = Query(...),
    period_id: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
):
    """Subjects the student archived or unenrolled (hidden unenrolled from My subjects)."""
    err = require_supabase()
    if err is not None:
        return err
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    allowed, _, bad = _can_view_student_data(authorization, sid)
    if not allowed:
        return bad
    student_uuid = db_supabase.profile_uuid_for_id_number(sid)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)
    try:
        subjects = db_supabase.list_archived_subjects_for_student(student_uuid, period_id)
        return {"subjects": subjects, "count": len(subjects)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.patch("/student/subjects/{subject_id}/enrollment")
async def patch_student_subject_enrollment_endpoint(
    subject_id: str,
    body: dict = Body(...),
    authorization: str | None = Header(default=None),
):
    """Archive or unenroll a student from a subject (updates Supabase enrollments)."""
    err = require_supabase()
    if err is not None:
        return err
    payload = body if isinstance(body, dict) else {}
    student_id_number, bad = resolve_student_id_number_or_403(payload, authorization)
    if bad is not None:
        return bad
    action = str(payload.get("action") or "").strip().lower()
    if not student_id_number:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    if action not in ("archive", "unenroll"):
        return JSONResponse(
            {"error": "action must be archive or unenroll"},
            status_code=400,
        )
    student_uuid = db_supabase.profile_uuid_for_id_number(student_id_number)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)
    period_id = payload.get("grading_period_id")
    if not period_id:
        cur = db_supabase.get_current_grading_period() or {}
        period_id = cur.get("id")
    new_status = (
        db_supabase.ENROLLMENT_STATUS_ARCHIVED
        if action == "archive"
        else db_supabase.ENROLLMENT_STATUS_UNENROLLED
    )
    try:
        row = db_supabase.get_student_enrollment_row(
            student_uuid, str(subject_id), period_id
        )
        if not row:
            return JSONResponse({"error": "Enrollment not found"}, status_code=404)
        updated = db_supabase.update_student_enrollment_status(
            student_uuid, str(subject_id), period_id, new_status
        )
        if not updated:
            return JSONResponse({"error": "Could not update enrollment"}, status_code=502)
        return {
            "ok": True,
            "action": action,
            "enrollment_status": new_status,
            "enrollment": updated,
        }
    except ValueError as ve:
        return JSONResponse({"error": str(ve)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _attach_subject_aggregate_counts(subjects: list[dict]) -> None:
    pub_counts = db_supabase.count_published_lessons_by_subject()
    total_counts = db_supabase.count_lessons_by_subject()
    teacher_counts = db_supabase.count_teachers_by_subject()
    for s in subjects:
        key = str(s.get("id") or "")
        s["published_lesson_count"] = pub_counts.get(key, 0)
        s["total_lesson_count"] = total_counts.get(key, 0)
        s["teacher_count"] = teacher_counts.get(key, 0)


def _subject_response(row: dict) -> dict:
    sid = row.get("id")
    return {
        "id": str(sid) if sid is not None else None,
        "name": row.get("name") or "",
        "description": row.get("description") or "",
        "color": row.get("color") or "",
        "join_code": row.get("join_code"),
        "created_by_teacher_id_number": row.get("created_by_teacher_id_number"),
        "deped_category": row.get("deped_category") or "academic_standard",
    }


@app.get("/subjects")
def list_subjects_endpoint(
    owner_teacher_id_number: str | None = Query(default=None),
):
    """List subjects with lesson counts. Filter by owner when owner_teacher_id_number is set."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        owner = (owner_teacher_id_number or "").strip()
        if owner:
            subjects = db_supabase.list_subjects_for_teacher_owner(owner)
        else:
            subjects = db_supabase.list_subjects()
        _attach_subject_aggregate_counts(subjects)
        return {"subjects": subjects, "count": len(subjects)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/subjects")
async def create_subject_endpoint(body: dict, authorization: str | None = Header(default=None)):
    """Create a new subject row (used by Teacher 'Add Subject' modal, and
    Admin's global Add Subject modal)."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    if role not in ("teacher", "admin"):
        return JSONResponse({"error": "Only teachers or admins can create subjects."}, status_code=403)

    name = (body.get("name") or "").strip()
    if not name:
        return JSONResponse({"error": "Subject name is required."}, status_code=400)
    description = body.get("description")
    color = body.get("color")
    # A teacher can only create a subject owned by themselves — the caller's
    # own token identity, not whatever the request body claims. Admin may
    # leave it unowned (school-wide) or set any teacher explicitly.
    if role == "teacher":
        created_by = caller_idn
    else:
        created_by = (body.get("created_by_teacher_id_number") or body.get("teacher_id_number") or "").strip() or None
    deped_category = body.get("deped_category")
    try:
        row = db_supabase.create_subject(
            name=name,
            description=description,
            color=color,
            created_by_teacher_id_number=created_by,
            deped_category=deped_category,
        )
        out = _subject_response(row)
        out["published_lesson_count"] = 0
        out["total_lesson_count"] = 0
        out["teacher_count"] = 0
        return out
    except ValueError as ve:
        return JSONResponse({"error": str(ve)}, status_code=400)
    except Exception as e:
        msg = str(e)
        # Friendly error if the database uniqueness constraint is violated.
        if "duplicate" in msg.lower() or "unique" in msg.lower():
            return JSONResponse({"error": "A subject with this name already exists."}, status_code=409)
        return JSONResponse({"error": msg}, status_code=502)


@app.post("/subjects/join")
async def join_subject_with_code_endpoint(body: dict = Body(...), authorization: str | None = Header(default=None)):
    """Student joins a subject using a join code."""
    err = require_supabase()
    if err is not None:
        return err
    payload = body if isinstance(body, dict) else {}
    join_code = str(payload.get("join_code") or payload.get("code") or "").strip()
    student_id_number, bad = resolve_student_id_number_or_403(payload, authorization)
    if bad is not None:
        return bad
    if not join_code:
        return JSONResponse({"error": "join_code is required"}, status_code=400)
    if not student_id_number:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    student_uuid = db_supabase.profile_uuid_for_id_number(student_id_number)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)
    try:
        result = db_supabase.join_subject_by_code(student_uuid, join_code)
        subject = result.get("subject") or {}
        return {
            "ok": True,
            "subject": _subject_response(subject),
            "enrollment": result.get("enrollment"),
        }
    except ValueError as ve:
        msg = str(ve)
        if "already enrolled" in msg.lower():
            return JSONResponse({"error": msg}, status_code=409)
        if "invalid subject code" in msg.lower():
            return JSONResponse({"error": msg}, status_code=404)
        return JSONResponse({"error": msg}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/subjects/{subject_id}/people")
def subject_people_endpoint(subject_id: str, authorization: str | None = Header(default=None)):
    """People tab: the subject's teacher + classmates roster. Only the
    owning teacher, an enrolled student, or an admin may view it."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    if not caller_prof:
        return JSONResponse({"error": "Profile not found."}, status_code=404)

    subject = db_supabase.get_subject_row(subject_id)
    if not subject:
        return JSONResponse({"error": "Subject not found."}, status_code=404)

    if not _subject_people_access(subject, caller_prof, caller_idn):
        return JSONResponse({"error": "You don't have access to this class."}, status_code=403)

    try:
        students = db_supabase.list_students_enrolled_in_subject(subject_id)
        teacher_idn = str(subject.get("created_by_teacher_id_number") or "").strip()
        teacher_prof = db_supabase.get_profile_by_id_number(teacher_idn) if teacher_idn else None
        teacher = db_supabase.serialize_public_profile(teacher_prof) if teacher_prof else None
        return {"teacher": teacher, "students": students, "count": len(students)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _subject_people_access(subject: dict, caller_prof: dict | None, caller_idn: str) -> bool:
    """Shared access check: owning teacher, an enrolled student, or admin."""
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    if role == "admin":
        return True
    if role == "teacher":
        return str(subject.get("created_by_teacher_id_number") or "").strip() == caller_idn
    if role == "student":
        student_uuid = str((caller_prof or {}).get("id") or "")
        return db_supabase.student_enrolled_in_subject(student_uuid, subject.get("id"))
    return False


@app.get("/subjects/{subject_id}/announcements")
def list_subject_announcements_endpoint(subject_id: str, authorization: str | None = Header(default=None)):
    """Class Stream feed — same access as the People tab (owning teacher,
    enrolled student, or admin)."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    if not caller_prof:
        return JSONResponse({"error": "Profile not found."}, status_code=404)
    subject = db_supabase.get_subject_row(subject_id)
    if not subject:
        return JSONResponse({"error": "Subject not found."}, status_code=404)
    if not _subject_people_access(subject, caller_prof, caller_idn):
        return JSONResponse({"error": "You don't have access to this class."}, status_code=403)
    try:
        return {"announcements": db_supabase.list_subject_announcements(subject_id, viewer_id_number=caller_idn)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/subjects/{subject_id}/announcements")
async def create_subject_announcement_endpoint(
    subject_id: str,
    body: dict = Body(...),
    authorization: str | None = Header(default=None),
):
    """Only the subject's owning teacher may post to the Class Stream."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    subject = db_supabase.get_subject_row(subject_id)
    if not subject:
        return JSONResponse({"error": "Subject not found."}, status_code=404)
    if role != "teacher" or str(subject.get("created_by_teacher_id_number") or "").strip() != caller_idn:
        return JSONResponse({"error": "Only this subject's teacher can post announcements."}, status_code=403)
    text = str((body or {}).get("body") or "").strip()
    if not text:
        return JSONResponse({"error": "Announcement text is required."}, status_code=400)
    lesson_id = (body or {}).get("lesson_id") or None
    try:
        db_supabase.create_subject_announcement(subject_id, caller_idn, text, lesson_id=lesson_id)
        return {"announcements": db_supabase.list_subject_announcements(subject_id, viewer_id_number=caller_idn)}
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _announcement_access_or_403(subject_id: str, announcement_id: str, authorization: str | None):
    """Shared guard for comment/reaction endpoints: caller must be the subject's
    teacher, an enrolled student, or admin — same rule as viewing the Class Stream.
    Returns (caller_idn, announcement_row, error_response_or_None)."""
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return None, None, JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    if not caller_prof:
        return None, None, JSONResponse({"error": "Profile not found."}, status_code=404)
    subject = db_supabase.get_subject_row(subject_id)
    if not subject:
        return None, None, JSONResponse({"error": "Subject not found."}, status_code=404)
    if not _subject_people_access(subject, caller_prof, caller_idn):
        return None, None, JSONResponse({"error": "You don't have access to this class."}, status_code=403)
    ann = db_supabase.get_announcement_row(announcement_id)
    if not ann or str(ann.get("subject_id") or "") != str(subject.get("id") or ""):
        return None, None, JSONResponse({"error": "Announcement not found."}, status_code=404)
    return caller_idn, caller_prof, None


@app.post("/subjects/{subject_id}/announcements/{announcement_id}/comments")
async def create_announcement_comment_endpoint(
    subject_id: str,
    announcement_id: str,
    body: dict = Body(...),
    authorization: str | None = Header(default=None),
):
    """Anyone with Class Stream access (owning teacher, enrolled student, admin) can comment."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn, caller_prof, bad = _announcement_access_or_403(subject_id, announcement_id, authorization)
    if bad is not None:
        return bad
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    comment_role = "teacher" if role in ("teacher", "admin") else "student"
    text = str((body or {}).get("body") or "").strip()
    if not text:
        return JSONResponse({"error": "Comment text is required."}, status_code=400)
    try:
        db_supabase.create_announcement_comment(announcement_id, caller_idn, comment_role, text)
        return {"announcements": db_supabase.list_subject_announcements(subject_id, viewer_id_number=caller_idn)}
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/subjects/{subject_id}/announcements/{announcement_id}/react")
async def toggle_announcement_reaction_endpoint(
    subject_id: str,
    announcement_id: str,
    authorization: str | None = Header(default=None),
):
    """Toggle the signed-in user's like on an announcement."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn, _caller_prof, bad = _announcement_access_or_403(subject_id, announcement_id, authorization)
    if bad is not None:
        return bad
    try:
        return db_supabase.toggle_announcement_reaction(announcement_id, caller_idn)
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/subjects/{subject_id}/regenerate-code")
async def regenerate_subject_join_code_endpoint(
    subject_id: str,
    body: dict = Body(...),
    authorization: str | None = Header(default=None),
):
    """Teacher regenerates join code; existing enrollments are kept."""
    err = require_supabase()
    if err is not None:
        return err
    if not subject_id:
        return JSONResponse({"error": "subject_id is required."}, status_code=400)
    payload = body if isinstance(body, dict) else {}
    teacher_id_number = str(
        payload.get("teacher_id_number") or payload.get("created_by_teacher_id_number") or ""
    ).strip()
    if not teacher_id_number:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    allowed, bad = _can_view_teacher_data(authorization, teacher_id_number)
    if not allowed:
        return bad
    try:
        row = db_supabase.regenerate_subject_join_code(subject_id, teacher_id_number)
        return _subject_response(row)
    except LookupError:
        return JSONResponse({"error": "Subject not found."}, status_code=404)
    except PermissionError as pe:
        return JSONResponse({"error": str(pe)}, status_code=403)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.put("/subjects/{subject_id}")
async def update_subject_endpoint(subject_id: str, body: dict, authorization: str | None = Header(default=None)):
    """Update a subject row (name/description/color). Admin-only."""
    err = require_supabase()
    if err is not None:
        return err
    if not subject_id:
        return JSONResponse({"error": "subject_id is required."}, status_code=400)
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        row = db_supabase.update_subject(
            subject_id=subject_id,
            name=body.get("name"),
            description=body.get("description"),
            color=body.get("color"),
            deped_category=body.get("deped_category"),
        )
        sid = row.get("id") if isinstance(row, dict) else subject_id
        return _subject_response(row if isinstance(row, dict) else {
            "id": sid,
            "name": body.get("name"),
            "description": body.get("description"),
            "color": body.get("color"),
            "deped_category": body.get("deped_category"),
        })
    except ValueError as ve:
        return JSONResponse({"error": str(ve)}, status_code=400)
    except Exception as e:
        msg = str(e)
        if "duplicate" in msg.lower() or "unique" in msg.lower():
            return JSONResponse({"error": "Another subject already uses this name."}, status_code=409)
        return JSONResponse({"error": msg}, status_code=502)


_LESSON_FILE_MEDIA = {
    "pdf": "application/pdf",
    "ppt": "application/vnd.ms-powerpoint",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _resolve_lesson_file_path(
    lesson_id: str,
    storage_path: str | None = None,
    filename: str | None = None,
) -> Path | None:
    """Find lesson file on disk (uploads/lessons, DB path, or legacy temp_upload)."""
    backend_dir = Path(__file__).resolve().parent
    lid = str(lesson_id or "").strip()

    if lid:
        for ext in (".pdf", ".ppt", ".pptx", ".PDF", ".PPT", ".PPTX"):
            canonical = LESSON_UPLOADS_DIR / f"{lid}{ext.lower()}"
            if canonical.is_file():
                return canonical.resolve()

    raw_path = (str(storage_path).strip() if storage_path else "") or ""
    if raw_path:
        candidates = [
            Path(raw_path),
            backend_dir / raw_path,
            UPLOADS_DIR / raw_path,
            BASE_DIR / raw_path,
            BASE_DIR / "backend" / raw_path,
        ]
        for candidate in candidates:
            try:
                resolved = candidate.resolve()
                if resolved.is_file():
                    return resolved
            except OSError:
                continue

    if filename:
        safe = Path(filename).name.replace(" ", "_")
        legacy_name = f"temp_upload_{safe}"
        for base in (backend_dir, BASE_DIR, BASE_DIR / "backend", Path.cwd()):
            try:
                legacy = (base / legacy_name).resolve()
                if legacy.is_file():
                    return legacy
            except OSError:
                continue
    return None


def _pptx_media_sort_key(name: str) -> tuple:
    m = re.search(r"image(\d+)", name, re.I)
    return (int(m.group(1)) if m else 0, name.lower())


def _pptx_collect_shape_text(shape, text_parts: list[str]) -> None:
    """Recursively collect text from shapes (groups, tables, text frames)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                _pptx_collect_shape_text(child, text_parts)
            return
    except Exception:
        pass

    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    block = (cell.text or "").strip()
                    if block:
                        text_parts.append(block)
        except Exception:
            pass
        return

    if getattr(shape, "has_text_frame", False):
        block = (shape.text or "").strip()
        if block:
            text_parts.append(block)


def _pptx_slide_images(source: Path | bytes, max_slides: int = LESSON_VISION_MAX_SLIDES) -> list[tuple[str, bytes]]:
    """JPEG/PNG embedded in PPTX (common when slides are exported as pictures)."""
    images: list[tuple[str, bytes]] = []
    try:
        filelike: io.BytesIO | Path = io.BytesIO(source) if isinstance(source, bytes) else source
        with zipfile.ZipFile(filelike) as zf:
            names = [
                n
                for n in zf.namelist()
                if n.startswith("ppt/media/")
                and n.lower().endswith((".png", ".jpg", ".jpeg", ".webp"))
            ]
            names.sort(key=_pptx_media_sort_key)
            for name in names[:max_slides]:
                ext = Path(name).suffix.lower()
                mime = "image/jpeg" if ext in (".jpg", ".jpeg") else "image/png"
                images.append((mime, zf.read(name)))
    except Exception as e:
        print(f"_pptx_slide_images: {e}")
    return images


def _gemini_ocr_lesson_images(images: list[tuple[str, bytes]]) -> str:
    """Use Gemini vision to read text from slide/page images (picture-only decks)."""
    if not images or not API_KEY or not str(API_KEY).strip():
        return ""

    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        f"gemini-2.5-flash:generateContent?key={API_KEY}"
    )
    prompt = (
        "These images are consecutive slides from a lesson presentation. "
        "Extract ALL readable text in slide order. Use plain text only—keep headings, "
        "bullets, and numbering. Separate slides with a blank line. "
        "If a slide has no readable text, write [slide: no text]."
    )
    parts: list[dict] = [{"text": prompt}]
    for mime, data in images:
        parts.append(
            {
                "inline_data": {
                    "mime_type": mime,
                    "data": base64.standard_b64encode(data).decode("ascii"),
                }
            }
        )
    try:
        response = requests.post(
            url,
            json={"contents": [{"parts": parts}]},
            timeout=180,
        )
        result = response.json()
        if response.status_code != 200:
            print(f"_gemini_ocr_lesson_images status {response.status_code}: {result}")
            return ""
        return (gemini_text_from_result(result) or "").strip()
    except Exception as e:
        print(f"_gemini_ocr_lesson_images: {e}")
        return ""


def extract_lesson_text_from_file(file_path: Path, filename: str | None = None) -> str:
    """Pull plain text from PDF or PPTX for AI reviewer/quiz generation."""
    name = (filename or file_path.name or "").lower()
    text_parts: list[str] = []

    if name.endswith(".pdf"):
        try:
            reader = PdfReader(str(file_path))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        except Exception as e:
            print(f"extract_lesson_text_from_file pdf: {e}")
    elif name.endswith(".pptx"):
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            for slide in prs.slides:
                for shape in slide.shapes:
                    _pptx_collect_shape_text(shape, text_parts)
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if notes:
                        text_parts.append(notes)
        except Exception as e:
            print(f"extract_lesson_text_from_file pptx: {e}")
    elif name.endswith(".ppt"):
        print("extract_lesson_text_from_file: legacy .ppt not supported; use .pptx or PDF")

    combined = "\n".join(text_parts).strip()
    if combined:
        return combined[:LESSON_EXTRACT_MAX_CHARS]
    return ""


def extract_lesson_text_from_bytes(data: bytes, filename: str | None = None) -> str:
    """Pull plain text from PDF or PPTX bytes (same rules as extract_lesson_text_from_file)."""
    name = (filename or "").lower()
    text_parts: list[str] = []
    buf = io.BytesIO(data)

    if name.endswith(".pdf"):
        try:
            reader = PdfReader(buf)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        except Exception as e:
            print(f"extract_lesson_text_from_bytes pdf: {e}")
    elif name.endswith(".pptx"):
        try:
            from pptx import Presentation

            buf.seek(0)
            prs = Presentation(buf)
            for slide in prs.slides:
                for shape in slide.shapes:
                    _pptx_collect_shape_text(shape, text_parts)
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                    if notes:
                        text_parts.append(notes)
        except Exception as e:
            print(f"extract_lesson_text_from_bytes pptx: {e}")
    elif name.endswith(".ppt"):
        print("extract_lesson_text_from_bytes: legacy .ppt not supported; use .pptx or PDF")

    combined = "\n".join(text_parts).strip()
    if combined:
        return combined[:LESSON_EXTRACT_MAX_CHARS]
    return ""


def extract_lesson_text_with_vision(source: Path | bytes, filename: str | None = None) -> str:
    """OCR fallback for picture-only PPTX slides via Gemini vision."""
    name = (filename or (source.name if isinstance(source, Path) else "") or "").lower()
    if not name.endswith(".pptx"):
        return ""
    images = _pptx_slide_images(source)
    if not images:
        return ""
    label = filename or (str(source) if isinstance(source, Path) else "pptx-bytes")
    print(f"extract_lesson_text_with_vision: {len(images)} slide image(s) from {label}")
    ocr = _gemini_ocr_lesson_images(images)
    return ocr[:LESSON_EXTRACT_MAX_CHARS] if ocr else ""


def _decode_lesson_file_bytes(lesson: dict) -> bytes | None:
    """Lesson file stored in lessons.file_base64 (standard base64)."""
    fb = lesson.get("file_base64")
    if fb is None or not str(fb).strip():
        return None
    try:
        return base64.standard_b64decode(str(fb).strip())
    except Exception as e:
        print(f"_decode_lesson_file_bytes: {e}")
        return None


def lesson_text_for_ai(lesson: dict, *, allow_vision_fallback: bool = False) -> tuple[str, str | None]:
    """
    Text used for Gemini prompts. Uses extracted_text, else file bytes in DB, else disk file.
    """
    existing = str(lesson.get("extracted_text") or "").strip()
    if existing:
        return existing, None

    lid = str(lesson.get("id") or "")
    fn = (lesson.get("filename") or "").lower()
    file_blob = _decode_lesson_file_bytes(lesson)

    if file_blob:
        fresh = extract_lesson_text_from_bytes(file_blob, lesson.get("filename"))
        if not fresh.strip() and allow_vision_fallback:
            fresh = extract_lesson_text_with_vision(file_blob, lesson.get("filename"))
        if fresh.strip():
            try:
                db_supabase.update_lesson_extracted_text(lid, fresh)
            except Exception as e:
                print(f"lesson_text_for_ai cache extract: {e}")
            return fresh, None

    file_path = _resolve_lesson_file_path(
        lid,
        lesson.get("storage_path"),
        lesson.get("filename"),
    )

    if file_path:
        fresh = extract_lesson_text_from_file(file_path, lesson.get("filename"))
        if not fresh.strip() and allow_vision_fallback:
            fresh = extract_lesson_text_with_vision(file_path, lesson.get("filename"))
        if fresh.strip():
            try:
                db_supabase.update_lesson_extracted_text(lid, fresh)
            except Exception as e:
                print(f"lesson_text_for_ai cache extract: {e}")
            return fresh, None

    if not file_blob and not file_path:
        if fn.endswith((".pptx", ".ppt", ".pdf")):
            return "", (
                "Lesson file is missing on the server (not in database and not on disk). "
                "Re-upload the presentation from the subject page, or restore uploads/lessons from backup."
            )
        return "", "No lesson file found. Upload a PDF or PowerPoint (.pptx) first."

    if fn.endswith(".pptx") or fn.endswith(".ppt"):
        return "", (
            "No readable text in this PowerPoint. Slides may be pictures only—we tried reading them "
            "but found nothing usable. Re-save with editable text boxes, or export as PDF with selectable text."
        )
    if fn.endswith(".pdf"):
        return "", (
            "No text found in this PDF. Use a file with selectable/copyable text (not a scanned photo PDF)."
        )
    return "", "No text extracted from this file. Upload a PDF with selectable text or a PPTX with slide text."


def _lesson_file_response(
    lesson_id: str,
    lesson: dict,
    *,
    as_attachment: bool = False,
) -> FileResponse | JSONResponse | Response:
    """Stream lesson file: prefer bytes in DB, else disk (uploads/lessons or legacy path)."""
    raw = _decode_lesson_file_bytes(lesson)
    if raw is not None:
        ext = Path(lesson.get("filename") or "lesson.bin").suffix.lower().lstrip(".") or "bin"
        media_type = _LESSON_FILE_MEDIA.get(ext, "application/octet-stream")
        filename = lesson.get("filename") or f"lesson.{ext}"
        disp = "attachment" if as_attachment else "inline"
        return Response(
            content=raw,
            media_type=media_type,
            headers={"Content-Disposition": f'{disp}; filename="{filename}"'},
        )

    file_path = _resolve_lesson_file_path(
        str(lesson_id),
        lesson.get("storage_path"),
        lesson.get("filename"),
    )
    if not file_path:
        sp = (lesson.get("storage_path") or "").strip()
        hint = f" (storage_path: {sp})" if sp else ""
        return JSONResponse(
            {
                "error": (
                    "Original lesson file is missing on the server disk. "
                    "The teacher may need to re-upload this lesson, or restore the file from backup."
                    + hint
                ),
            },
            status_code=404,
        )
    try:
        LESSON_UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
        canonical = LESSON_UPLOADS_DIR / f"{lesson_id}{file_path.suffix.lower()}"
        if not canonical.is_file():
            shutil.copy2(file_path, canonical)
            db_supabase.update_lesson_storage_path(str(lesson_id), f"lessons/{canonical.name}")
            file_path = canonical
    except Exception as copy_err:
        print(f"lesson file canonicalize: {copy_err}")
    ext = file_path.suffix.lower().lstrip(".")
    media_type = _LESSON_FILE_MEDIA.get(ext, "application/octet-stream")
    filename = lesson.get("filename") or file_path.name
    disposition = "attachment" if as_attachment else "inline"
    return FileResponse(
        path=str(file_path),
        media_type=media_type,
        filename=filename,
        content_disposition_type=disposition,
    )


@app.get("/lessons/{lesson_id}/file")
def view_lesson_file(
    lesson_id: str,
    teacher_id_number: str = Query(...),
    download: bool = Query(False),
):
    """Stream the uploaded lesson file (teacher must own the lesson)."""
    err = require_supabase()
    if err is not None:
        return err
    tid = (teacher_id_number or "").strip()
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required."}, status_code=400)
    try:
        lesson = db_supabase.get_lesson_row(str(lesson_id))
        if not lesson:
            return JSONResponse({"error": "Lesson not found."}, status_code=404)
        owner = (lesson.get("teacher_id_number") or "").strip()
        if owner and owner != tid:
            return JSONResponse({"error": "You can only view your own lessons."}, status_code=403)
        return _lesson_file_response(str(lesson_id), lesson, as_attachment=download)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/lessons/{lesson_id}/file")
def view_student_lesson_file(
    lesson_id: str,
    student_id_number: str = Query(...),
    download: bool = Query(False),
    authorization: str | None = Header(default=None),
):
    """Stream a published lesson file for an enrolled student."""
    err = require_supabase()
    if err is not None:
        return err
    sid, auth_err = resolve_student_id_number_or_403(
        {"student_id_number": student_id_number},
        authorization,
    )
    if auth_err is not None:
        return auth_err
    try:
        lesson = db_supabase.get_lesson_row(str(lesson_id))
        if not lesson:
            return JSONResponse({"error": "Lesson not found."}, status_code=404)
        student_uuid = db_supabase.profile_uuid_for_id_number(sid)
        if not student_uuid:
            return JSONResponse({"error": "Student not found"}, status_code=404)
        if not db_supabase.student_can_view_published_lesson(student_uuid, lesson):
            return JSONResponse(
                {"error": "You do not have access to this lesson."},
                status_code=403,
            )
        return _lesson_file_response(str(lesson_id), lesson, as_attachment=download)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.delete("/lessons/{lesson_id}")
async def delete_lesson_endpoint(lesson_id: str, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    if not lesson_id:
        return JSONResponse({"error": "lesson_id is required."}, status_code=400)
    try:
        lesson = db_supabase.get_lesson_row(str(lesson_id))
        if not lesson:
            return JSONResponse({"error": "Lesson not found."}, status_code=404)
        allowed, bad = _can_view_teacher_data(authorization, str(lesson.get("teacher_id_number") or ""))
        if not allowed:
            return bad
        db_supabase.delete_lesson(str(lesson_id))
        return {"deleted_lesson_id": lesson_id}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.delete("/subjects/{subject_id}")
async def delete_subject_endpoint(subject_id: str, authorization: str | None = Header(default=None)):
    """Delete a subject. Any lesson referencing it will have its subject_id set
    to NULL first (so legacy lessons appear under 'Unassigned'). The owning
    teacher or an admin may delete it."""
    err = require_supabase()
    if err is not None:
        return err
    if not subject_id:
        return JSONResponse({"error": "subject_id is required."}, status_code=400)
    subject = db_supabase.get_subject_row(subject_id)
    if not subject:
        return JSONResponse({"error": "Subject not found."}, status_code=404)
    allowed, bad = _can_view_teacher_data(authorization, str(subject.get("created_by_teacher_id_number") or ""))
    if not allowed:
        return bad
    try:
        db_supabase.delete_subject(subject_id)
        return {"deleted_subject_id": subject_id}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/lesson/subject")
async def set_lesson_subject(body: dict, authorization: str | None = Header(default=None)):
    """Teacher edit: assign or change the subject of an existing lesson."""
    err = require_supabase()
    if err is not None:
        return err
    lesson_id = body.get("lesson_id") or body.get("file_id")
    if not lesson_id:
        return JSONResponse({"error": "lesson_id (or file_id) is required."}, status_code=400)
    subject_id = body.get("subject_id")
    if subject_id is not None:
        subject_id = str(subject_id) or None
    try:
        lesson = db_supabase.get_lesson_row(str(lesson_id))
        if not lesson:
            return JSONResponse({"error": "Lesson not found."}, status_code=404)
        allowed, bad = _can_view_teacher_data(authorization, str(lesson.get("teacher_id_number") or ""))
        if not allowed:
            return bad
        db_supabase.update_lesson_subject(str(lesson_id), subject_id)
        return {"lesson_id": lesson_id, "subject_id": subject_id}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/lesson")
def get_student_lesson():
    err = require_supabase()
    if err is not None:
        return err
    try:
        bundle = db_supabase.get_published_lesson_with_content()
        if not bundle:
            return JSONResponse(
                {"error": "No lesson has been published yet. Ask your teacher to publish a lesson."},
                status_code=404,
            )
        meta, gen = bundle
        lesson_id = meta["id"]
        reviewer = gen.get("reviewer")
        if isinstance(reviewer, list):
            reviewer_str = "\n\n".join(str(x).strip() for x in reviewer if str(x).strip())
        elif isinstance(reviewer, str):
            reviewer_str = reviewer
        else:
            reviewer_str = ""

        quiz = gen.get("quiz") or []
        if not isinstance(quiz, list):
            if isinstance(quiz, dict):
                # Convert single question object to array
                quiz = [quiz]
            else:
                quiz = []

        activities = gen.get("activities") or []
        print(f"[DEBUG] Raw activities from database: {activities}")
        print(f"[DEBUG] Type of raw activities: {type(activities)}")
        if not isinstance(activities, list):
            if isinstance(activities, str):
                # Convert single string to array
                activities = [activities]
                print(f"[DEBUG] Converted string to array: {activities}")
            else:
                activities = []
                print(f"[DEBUG] Converted non-list to empty array")
        else:
            print(f"[DEBUG] Activities already an array: {activities}")

        result = {
            "file_id": lesson_id,
            "filename": meta.get("filename", ""),
            "reviewer": reviewer_str,
            "quiz": quiz,
            "activities": activities,
        }
        print(f"[DEBUG] /student/lesson returning activities: {activities}")
        return result
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/leaderboard")
def get_student_leaderboard(limit: int = Query(50, ge=1, le=200)):
    """Live rankings from quiz_attempts (students only)."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        return db_supabase.get_learniq_leaderboard(limit=limit)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/dashboard-stats")
def get_student_dashboard_stats(authorization: str | None = Header(default=None)):
    """LearnIQ dashboard: quiz totals, rank, preview (Bearer token)."""
    err = require_supabase()
    if err is not None:
        return err
    sid = student_id_number_from_authorization(authorization)
    if not sid:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        return db_supabase.get_student_learniq_dashboard_stats(sid)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/learning-iq")
def get_student_learning_iq_endpoint(authorization: str | None = Header(default=None)):
    """LearnIQ Learning IQ score from quiz attempts and learning events (Bearer token)."""
    err = require_supabase()
    if err is not None:
        return err
    sid = student_id_number_from_authorization(authorization)
    if not sid:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        return db_supabase.compute_student_learning_iq(sid)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/upload-lesson")
async def upload_lesson_json(body: dict, authorization: str | None = Header(default=None)):
    """Create a lesson row from JSON (same data as your Flask sample, without a file upload)."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    role = str((caller_prof or {}).get("role") or "").strip().lower()
    if role not in ("teacher", "admin"):
        return JSONResponse({"error": "Only teachers or admins can upload lessons."}, status_code=403)
    try:
        filename = (body.get("filename") or "").strip()
        if not filename:
            return JSONResponse({"error": "filename is required."}, status_code=400)
        text = body.get("extracted_text") or body.get("text") or ""
        ft = (body.get("file_type") or Path(filename).suffix.lstrip(".") or "unknown").strip()
        # Teachers can only upload under their own name — the token identity,
        # not whatever teacher_id_number the request body claims.
        tid = caller_idn if role == "teacher" else (body.get("teacher_id_number") or body.get("teacher_id"))
        if tid is not None:
            tid = str(tid)
        subject_id = body.get("subject_id")
        if subject_id is not None:
            subject_id = str(subject_id) or None
        lesson = db_supabase.insert_lesson(
            filename=filename,
            file_type=ft,
            extracted_text=str(text)[:3000],
            storage_path=None,
            teacher_id_number=tid,
            subject_id=subject_id,
        )
        return lesson
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/upload-file")
async def upload_file(
    file: UploadFile = File(...),
    teacher_id_number: str = Form(...),
    subject_id: str | None = Form(default=None),
    authorization: str | None = Header(default=None),
):
    print("UPLOAD CALLED")
    print("Filename:", file.filename)
    print("Teacher ID Number:", teacher_id_number)
    print("Subject ID:", subject_id)

    err = require_supabase()
    if err is not None:
        return err

    allowed, _ = _can_view_teacher_data(authorization, teacher_id_number)
    if not allowed:
        return JSONResponse({"error": "You can only upload lessons under your own teacher account."}, status_code=403)

    # Recreate dirs if someone deleted `uploads/` while the API is running (mount + any legacy paths).
    UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
    LESSON_UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

    if not file.filename or not file.filename.lower().endswith((".pdf", ".ppt", ".pptx")):
        return JSONResponse(
            {"error": "Only PDF, PPT, and PPTX files are allowed."},
            status_code=400,
        )

    ext = Path(file.filename).suffix.lower()
    file_type = ext.lstrip(".") or "unknown"

    raw = await file.read()
    if len(raw) > LESSON_UPLOAD_MAX_BYTES:
        mb = LESSON_UPLOAD_MAX_BYTES // (1024 * 1024)
        return JSONResponse(
            {"error": f"File too large (max {mb} MB). Export a smaller PDF/PPTX or split the deck."},
            status_code=413,
        )

    try:
        text = await asyncio.wait_for(
            asyncio.to_thread(extract_lesson_text_from_bytes, raw, file.filename),
            timeout=LESSON_TEXT_EXTRACT_TIMEOUT_SEC,
        )
    except asyncio.TimeoutError:
        print(
            "upload-file: extract_lesson_text_from_bytes timed out; "
            f"saving with empty extracted_text ({file.filename!r})"
        )
        text = ""

    clean_subject_id = (str(subject_id).strip() or None) if subject_id else None
    lesson_uuid = str(uuid.uuid4())
    rel_storage = f"lessons/{lesson_uuid}{ext}"
    dest_path = LESSON_UPLOADS_DIR / f"{lesson_uuid}{ext}"

    try:
        print("CALLING INSERT LESSON (disk storage, metadata only)...")
        dest_path.write_bytes(raw)
        try:
            lesson = db_supabase.insert_lesson(
                filename=file.filename,
                file_type=file_type,
                extracted_text=text,
                storage_path=rel_storage,
                teacher_id_number=teacher_id_number,
                subject_id=clean_subject_id,
                file_base64=None,
                lesson_id=lesson_uuid,
            )
        except Exception:
            try:
                dest_path.unlink(missing_ok=True)
            except OSError as unlink_err:
                print(f"upload-file: rollback unlink failed: {unlink_err}")
            raise
        lid = str(lesson["id"])
        print(
            f"UPLOAD SUCCESS: file_id={lid}, filename={file.filename}, "
            f"subject_id={clean_subject_id}, storage_path={rel_storage} (local disk)"
        )
        return {"file_id": lid, "filename": file.filename, "subject_id": clean_subject_id}
    except Exception as e:
        import traceback

        print(f"UPLOAD FAILED: {e}")
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/generate-reviewer")
async def generate_reviewer(body: dict, authorization: str | None = Header(default=None)):
    print("AI GENERATION REQUEST RECEIVED: /generate-reviewer")
    print("REQUEST PAYLOAD:", body)
    auth_err = _require_signed_in(authorization)
    if auth_err is not None:
        return auth_err
    key_err = require_gemini_key()
    if key_err is not None:
        print("AI GENERATION ERROR (gemini key):", key_err.body if hasattr(key_err, "body") else key_err)
        return key_err
    db_err = require_supabase()
    if db_err is not None:
        print("AI GENERATION ERROR (supabase):", db_err.body if hasattr(db_err, "body") else db_err)
        return db_err

    file_id = body.get("file_id")
    lesson = db_supabase.get_lesson_row(str(file_id)) if file_id else None
    if not lesson:
        return JSONResponse({"error": "File not found"}, status_code=404)

    if not body.get("skip_cooldown"):
        cd_err = check_ai_generation_cooldown("reviewer", str(file_id))
        if cd_err is not None:
            return cd_err

    text, text_err = lesson_text_for_ai(lesson, allow_vision_fallback=True)
    if text_err:
        print("AI GENERATION ERROR: empty extracted_text")
        return JSONResponse({"error": text_err}, status_code=400)

    source = str(text).replace("\r\n", "\n").replace("\r", "\n")
    if len(source) > REVIEWER_SOURCE_MAX_CHARS:
        source = source[:REVIEWER_SOURCE_MAX_CHARS] + "\n\n[…excerpt truncated for length…]"
    prompt = REVIEWER_PROMPT_TEMPLATE.format(source=source)
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}

    print("SENDING TO AI API (reviewer)")
    response = requests.post(url, json=payload, timeout=120)
    result = response.json()
    print("AI RAW RESPONSE STATUS (reviewer):", response.status_code)
    try:
        print("AI RAW RESPONSE BODY (reviewer):", result)
    except Exception as e:
        print("AI RAW RESPONSE PRINT ERROR (reviewer):", str(e))

    if response.status_code != 200:
        return JSONResponse({"error": result}, status_code=502)

    reviewer_text = strip_outer_markdown_code_fence(gemini_text_from_result(result))
    if not reviewer_text:
        return JSONResponse({"error": "AI returned no text. Try again."}, status_code=502)

    try:
        db_supabase.set_reviewer(str(file_id), reviewer_text)
    except Exception as e:
        print("AI GENERATION ERROR (db write reviewer):", str(e))
        return JSONResponse({"error": str(e)}, status_code=502)

    if not body.get("skip_cooldown"):
        start_ai_generation_cooldown("reviewer", str(file_id))
    return {"reviewer": reviewer_text}


@app.post("/generate-question")
async def generate_question(body: dict, authorization: str | None = Header(default=None)):
    print("AI GENERATION REQUEST RECEIVED: /generate-question")
    print("REQUEST PAYLOAD:", body)
    auth_err = _require_signed_in(authorization)
    if auth_err is not None:
        return auth_err
    key_err = require_gemini_key()
    if key_err is not None:
        print("AI GENERATION ERROR (gemini key):", key_err.body if hasattr(key_err, "body") else key_err)
        return key_err
    db_err = require_supabase()
    if db_err is not None:
        print("AI GENERATION ERROR (supabase):", db_err.body if hasattr(db_err, "body") else db_err)
        return db_err

    file_id = body.get("file_id")
    lesson = db_supabase.get_lesson_row(str(file_id)) if file_id else None
    if not lesson:
        return JSONResponse({"error": "File not found"}, status_code=404)

    if not body.get("skip_cooldown"):
        cd_err = check_ai_generation_cooldown("quiz", str(file_id))
        if cd_err is not None:
            return cd_err

    text, text_err = lesson_text_for_ai(lesson, allow_vision_fallback=True)
    if text_err:
        print("AI GENERATION ERROR: empty extracted_text")
        return JSONResponse({"error": text_err}, status_code=400)

    quiz_count = body.get("quiz_count", 1)
    if not isinstance(quiz_count, int) or quiz_count < 1:
        quiz_count = 1
    difficulty = (body.get("difficulty") or "").strip().lower()
    if difficulty not in ("easy", "medium", "hard", ""):
        difficulty = ""
    
    diff_line = f"Difficulty: {difficulty}.\n" if difficulty else ""
    prompt = (
        "You are an educational content generator.\n"
        f"{diff_line}"
        f"Create exactly {quiz_count} multiple-choice questions based ONLY on the lesson text.\n"
        "Each question must have exactly 4 choices (A-D).\n"
        "Return STRICT VALID JSON ONLY (no markdown, no explanations, no code fences) with this schema:\n"
        '{ "questions": [ { "question": "...", "choices": ["A. ...","B. ...","C. ...","D. ..."], "answer": "B" } ] }\n'
        "The answer MUST be a single letter A, B, C, or D.\n\n"
        "LESSON TEXT:\n"
        f"{text}"
    )
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}

    print("SENDING TO AI API (quiz)")
    response = requests.post(url, json=payload, timeout=120)
    result = response.json()
    print("AI RAW RESPONSE STATUS (quiz):", response.status_code)
    try:
        print("AI RAW RESPONSE BODY (quiz):", result)
    except Exception as e:
        print("AI RAW RESPONSE PRINT ERROR (quiz):", str(e))

    if response.status_code != 200:
        return JSONResponse({"error": friendly_ai_error(result)}, status_code=502)

    raw_output = gemini_text_from_result(result)
    if not raw_output:
        return JSONResponse({"error": "AI returned no text. Try again."}, status_code=502)

    try:
        parsed = parse_model_json(raw_output)
        questions_data = normalize_quiz_questions(parsed, desired_count=quiz_count)
    except (json.JSONDecodeError, ValueError) as e:
        print("AI GENERATION ERROR: failed to parse questions JSON")
        return JSONResponse(
            {"error": "Failed to parse quiz questions. Please retry."},
            status_code=502,
        )

    try:
        # Clear existing quiz and add all new questions
        db_supabase.set_quiz(str(file_id), [])
        for question in questions_data:
            db_supabase.append_quiz_question(str(file_id), question)
    except Exception as e:
        print("AI GENERATION ERROR (db write quiz):", str(e))
        return JSONResponse({"error": str(e)}, status_code=502)

    if not body.get("skip_cooldown"):
        start_ai_generation_cooldown("quiz", str(file_id))
    return {"questions": questions_data, "count": len(questions_data)}


def normalize_battle_questions(questions: object) -> list[dict[str, str]]:
    if not isinstance(questions, list):
        return []
    seen_answers: set[str] = set()
    normalized: list[dict[str, str]] = []
    for entry in questions:
        if not isinstance(entry, dict):
            continue
        question_text = str(entry.get("question") or "").strip()
        if not question_text:
            continue
        if len(question_text) > 220:
            question_text = question_text[:217].rstrip() + "..."
        answer = re.sub(r"[^a-z]", "", str(entry.get("answer") or "").strip().lower())
        if len(answer) < 4 or len(answer) > 10:
            continue
        if answer in seen_answers:
            continue
        seen_answers.add(answer)
        meaning = str(entry.get("meaning") or "").strip()
        if len(meaning) > 140:
            meaning = meaning[:137].rstrip() + "..."
        normalized.append({"question": question_text, "answer": answer, "meaning": meaning})
    return normalized[:12]


@app.post("/generate-battle-questions")
async def generate_battle_questions(body: dict, authorization: str | None = Header(default=None)):
    print("AI GENERATION REQUEST RECEIVED: /generate-battle-questions")
    auth_err = _require_signed_in(authorization)
    if auth_err is not None:
        return auth_err
    key_err = require_gemini_key()
    if key_err is not None:
        return key_err
    db_err = require_supabase()
    if db_err is not None:
        return db_err

    file_id = body.get("file_id")
    lesson = db_supabase.get_lesson_row(str(file_id)) if file_id else None
    if not lesson:
        return JSONResponse({"error": "File not found"}, status_code=404)

    if not body.get("skip_cooldown"):
        cd_err = check_ai_generation_cooldown("battle_questions", str(file_id))
        if cd_err is not None:
            return cd_err

    text, text_err = lesson_text_for_ai(lesson, allow_vision_fallback=True)
    if text_err:
        return JSONResponse({"error": text_err}, status_code=400)

    prompt = (
        "You are an educational game designer creating a quiz-battle game. The student reads a "
        "question, then must spell the answer using letter tiles to attack an opponent. A wrong "
        "answer lets the opponent attack the student instead, so questions must be answerable "
        "from the lesson text alone.\n"
        "TASK: From the lesson text below, write 10 to 12 short questions, each with exactly ONE "
        "single-word answer.\n"
        "Rules:\n"
        "- Each question must be answerable with exactly ONE word — no phrases, no multi-word answers.\n"
        "- The answer word must be letters only (no spaces/punctuation/numbers), 4 to 10 letters long.\n"
        "- Base every question only on facts, terms, or concepts that actually appear in the lesson text.\n"
        "- Keep each question under 20 words and unambiguous.\n"
        "- Also include a one-sentence meaning/explanation of the answer term.\n"
        "Return STRICT VALID JSON ONLY (no markdown, no explanations, no code fences) with this schema:\n"
        '{ "questions": [ { "question": "...", "answer": "...", "meaning": "..." } ] }\n\n'
        "LESSON TEXT:\n"
        f"{text}"
    )
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}

    print("SENDING TO AI API (battle_questions)")
    response = requests.post(url, json=payload, timeout=120)
    result = response.json()
    print("AI RAW RESPONSE STATUS (battle_questions):", response.status_code)

    if response.status_code != 200:
        return JSONResponse({"error": friendly_ai_error(result)}, status_code=502)

    raw_output = gemini_text_from_result(result)
    if not raw_output:
        return JSONResponse({"error": "AI returned no text. Try again."}, status_code=502)

    try:
        parsed = parse_model_json(raw_output)
        questions = normalize_battle_questions(parsed.get("questions") if isinstance(parsed, dict) else None)
    except (json.JSONDecodeError, ValueError):
        return JSONResponse({"error": "Failed to parse battle questions. Please retry."}, status_code=502)

    if len(questions) < 5:
        return JSONResponse({"error": "AI could not build enough usable questions from this lesson."}, status_code=502)

    try:
        db_supabase.set_battle_questions(str(file_id), questions)
    except Exception as e:
        print("AI GENERATION ERROR (db write battle_questions):", str(e))
        return JSONResponse({"error": str(e)}, status_code=502)

    if not body.get("skip_cooldown"):
        start_ai_generation_cooldown("battle_questions", str(file_id))
    return {"questions": questions}


@app.post("/generate-activities")
async def generate_activities(body: dict, authorization: str | None = Header(default=None)):
    print(f"[DEBUG] /generate-activities called with body: {body}")
    auth_err = _require_signed_in(authorization)
    if auth_err is not None:
        return auth_err

    key_err = require_gemini_key()
    if key_err is not None:
        print(f"[DEBUG] Gemini key error: {key_err}")
        return key_err
    db_err = require_supabase()
    if db_err is not None:
        print(f"[DEBUG] Supabase error: {db_err}")
        return db_err

    file_id = body.get("file_id")
    print(f"[DEBUG] file_id: {file_id}")
    
    lesson = db_supabase.get_lesson_row(str(file_id)) if file_id else None
    if not lesson:
        print(f"[DEBUG] Lesson not found for file_id: {file_id}")
        return JSONResponse({"error": "File not found"}, status_code=404)

    if not body.get("skip_cooldown"):
        cd_err = check_ai_generation_cooldown("activity", str(file_id))
        if cd_err is not None:
            return cd_err

    text, text_err = lesson_text_for_ai(lesson, allow_vision_fallback=True)
    print(f"[DEBUG] Extracted text length: {len(text)}")
    if text_err:
        print("[DEBUG] No extracted text found")
        return JSONResponse({"error": text_err}, status_code=400)

    activity_type = (body.get("activity_type") or "essay").strip().lower()
    count = body.get("count", 5)
    try:
        count = int(count)
    except Exception:
        count = 5
    if count not in (5, 10, 15):
        count = 5

    allowed = {"essay", "flashcards"}
    if activity_type not in allowed:
        activity_type = "essay"

    schema_by_type = {
        "essay": '{ "activities": [ { "question": "Essay prompt...", "answer": "Sample answer or key points the response should cover (2-4 sentences)." } ] }',
        "flashcards": '{ "cards": [ { "front": "Term or question (short)", "back": "Definition or answer (concise)" } ] }',
    }

    type_instructions = {
        "essay": (
            "Write open-ended essay prompts that require the student to explain, analyze, "
            "compare, or evaluate ideas from the lesson. Each prompt should be 1-2 sentences. "
            "The 'answer' field should be a concise sample answer or list of key points (no markdown)."
        ),
        "flashcards": (
            "Create study flashcards. The 'front' is a short term, concept, or question "
            "(max ~8 words). The 'back' is a concise definition or answer (1-2 sentences). "
            "Cover the most important concepts from the lesson."
        ),
    }

    prompt = (
        "You are an educational content generator.\n"
        f"Create exactly {count} items for activity type: {activity_type}.\n"
        f"{type_instructions[activity_type]}\n"
        "Return STRICT VALID JSON ONLY (no markdown, no explanations, no code fences).\n"
        f"Schema:\n{schema_by_type[activity_type]}\n\n"
        "LESSON TEXT:\n"
        f"{text}"
    )
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}

    response = requests.post(url, json=payload, timeout=120)
    result = response.json()
    print(f"[DEBUG] Gemini API response status: {response.status_code}")

    if response.status_code != 200:
        print(f"[DEBUG] Gemini API error: {result}")
        return JSONResponse({"error": friendly_ai_error(result)}, status_code=502)

    raw_output = gemini_text_from_result(result)
    print(f"[DEBUG] Raw output from Gemini: {raw_output}")
    if not raw_output:
        print(f"[DEBUG] Gemini returned no text")
        return JSONResponse({"error": "AI returned no text. Try again."}, status_code=502)

    try:
        activity_data = parse_model_json(raw_output)
        if not isinstance(activity_data, dict):
            raise ValueError("Invalid activities JSON format")
    except (json.JSONDecodeError, ValueError) as e:
        print(f"[DEBUG] JSON parsing error: {e}")
        return JSONResponse({"error": "Failed to generate activities. Please retry."}, status_code=502)

    # Normalize to DB storage format: list of dicts, or a single deck object for flashcards.
    normalized_activities = []
    if activity_type == "flashcards":
        cards = activity_data.get("cards") if isinstance(activity_data.get("cards"), list) else []
        cards = [
            {"front": str(c.get("front") or "").strip(), "back": str(c.get("back") or "").strip()}
            for c in cards
            if isinstance(c, dict) and c.get("front") and c.get("back")
        ]
        if not cards:
            return JSONResponse({"error": "Failed to generate flashcards. Please retry."}, status_code=502)
        normalized_activities = [{"activity_type": "flashcards", "cards": cards}]
    else:
        acts = activity_data.get("activities") if isinstance(activity_data.get("activities"), list) else []
        for a in acts:
            if not isinstance(a, dict):
                continue
            q = str(a.get("question") or "").strip()
            ans = a.get("answer")
            if not q:
                continue
            normalized_activities.append({"activity_type": activity_type, "question": q, "answer": ans})
        if not normalized_activities:
            return JSONResponse({"error": "Failed to generate activities. Please retry."}, status_code=502)

    try:
        # Replace activities cleanly (stable regeneration)
        print(f"[DEBUG] Saving activities to database for file_id: {file_id}")
        db_supabase.set_activities(str(file_id), normalized_activities)
        print(f"[DEBUG] Activities saved successfully")
    except Exception as e:
        print(f"[DEBUG] Database save error: {e}")
        return JSONResponse({"error": str(e)}, status_code=502)

    if not body.get("skip_cooldown"):
        start_ai_generation_cooldown("activity", str(file_id))
    print(f"[DEBUG] Returning activities (normalized): {normalized_activities}")
    return {"activities": normalized_activities, "total_activities": len(normalized_activities)}


@app.post("/save-ai-content")
async def save_ai_content(body: dict, authorization: str | None = Header(default=None)):
    """Save reviewer, quiz, and/or activities for a lesson (manual or external tools)."""
    err = require_supabase()
    if err is not None:
        return err
    auth_err = _require_signed_in(authorization)
    if auth_err is not None:
        return auth_err
    lesson_id = body.get("lesson_id") or body.get("file_id")
    if not lesson_id:
        return JSONResponse({"error": "lesson_id or file_id is required."}, status_code=400)
    lesson_id = str(lesson_id)
    if not db_supabase.get_lesson_row(lesson_id):
        return JSONResponse({"error": "Lesson not found."}, status_code=404)
    try:
        if "reviewer" in body and body["reviewer"] is not None:
            db_supabase.set_reviewer(lesson_id, body["reviewer"])
        if "activities" in body and body["activities"] is not None:
            db_supabase.set_activities(lesson_id, body["activities"])
        if "quiz" in body and body["quiz"] is not None:
            q = body["quiz"]
            if not isinstance(q, list):
                return JSONResponse({"error": "quiz must be a JSON array."}, status_code=400)
            db_supabase.set_quiz(lesson_id, q)
        gen = db_supabase.get_content_row(lesson_id) or {}
        lesson = db_supabase.get_lesson_row(lesson_id) or {}
        return {
            "file_id": lesson_id,
            "filename": lesson.get("filename", ""),
            "reviewer": gen.get("reviewer"),
            "quiz": gen.get("quiz") or [],
            "activities": gen.get("activities"),
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/get-content/{file_id}")
def get_content(file_id: str):
    err = require_supabase()
    if err is not None:
        return err
    try:
        lesson = db_supabase.get_lesson_row(file_id)
        if not lesson:
            return JSONResponse({"error": "Content not found"}, status_code=404)
        gen = db_supabase.get_content_row(file_id) or {}
        return {
            "file_id": file_id,
            "filename": lesson.get("filename", ""),
            "reviewer": gen.get("reviewer"),
            "quiz": gen.get("quiz") or [],
            "activities": gen.get("activities"),
            "battle_questions": gen.get("battle_questions") or [],
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


# --- Quiz, attendance, journals ---


@app.get("/student/learning-history")
def student_learning_history_endpoint(
    student_id_number: str = Query(...),
    authorization: str | None = Header(default=None),
):
    """Quiz / reviewer / activity history from database (History page)."""
    err = require_supabase()
    if err is not None:
        return err
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    allowed, _, bad = _can_view_student_data(authorization, sid)
    if not allowed:
        return bad
    try:
        data = db_supabase.get_student_learning_history(sid)
        return {
            "student_id_number": sid,
            "quiz": data.get("quiz") or [],
            "reviewer": data.get("reviewer") or [],
            "activity": data.get("activity") or [],
            "battle": data.get("battle") or [],
            "counts": {
                "quiz": len(data.get("quiz") or []),
                "reviewer": len(data.get("reviewer") or []),
                "activity": len(data.get("activity") or []),
                "battle": len(data.get("battle") or []),
            },
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/student/learning-history")
async def student_learning_history_post_endpoint(
    body: dict = Body(...),
    authorization: str | None = Header(default=None),
):
    """Record reviewer or activity history event."""
    err = require_supabase()
    if err is not None:
        return err
    payload = body if isinstance(body, dict) else {}
    sid, bad = resolve_student_id_number_or_403(payload, authorization)
    if bad is not None:
        return bad
    event_type = str(payload.get("event_type") or payload.get("type") or "").strip().lower()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    if event_type not in ("reviewer", "activity", "battle"):
        return JSONResponse(
            {"error": "event_type must be reviewer, activity, or battle"},
            status_code=400,
        )
    try:
        row = db_supabase.insert_student_learning_event(sid, event_type, payload)
        return {"ok": True, "event": row}
    except ValueError as ve:
        return JSONResponse({"error": str(ve)}, status_code=400)
    except Exception as e:
        msg = str(e)
        if "student_learning_events" in msg and "does not exist" in msg.lower():
            return JSONResponse(
                {
                    "error": "Run backend/migrations/student_learning_events.sql in Supabase first.",
                },
                status_code=503,
            )
        return JSONResponse({"error": msg}, status_code=502)


@app.post("/quiz-attempt")
async def quiz_attempt(body: dict, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    try:
        lesson_id = body.get("lesson_id") or body.get("file_id")
        if not lesson_id:
            return JSONResponse({"error": "lesson_id (or file_id) is required."}, status_code=400)
        student_id_number, bad = resolve_student_id_number_or_403(body, authorization)
        if bad is not None:
            return bad
        row = db_supabase.insert_quiz_attempt(
            str(lesson_id),
            score=int(body.get("score", 0)),
            total_questions=int(body.get("total_questions", 0)),
            answers=body.get("answers"),
            student_id_number=student_id_number,
        )
        return row
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _parse_capture_timestamp_iso(raw: str) -> datetime:
    s = (raw or "").strip()
    if not s:
        raise ValueError("capture_timestamp is required.")
    if s.endswith("Z"):
        s = s[:-1] + "+00:00"
    dt = datetime.fromisoformat(s)
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    return dt.astimezone(timezone.utc)


def _reverse_geocode_location(lat: float, lon: float) -> str | None:
    try:
        res = requests.get(
            "https://nominatim.openstreetmap.org/reverse",
            params={"lat": lat, "lon": lon, "format": "json"},
            headers={"User-Agent": "LearnIQTrack/1.0 (immersion attendance)"},
            timeout=10,
        )
        if res.status_code != 200:
            return None
        data = res.json()
        return (data.get("display_name") or "").strip() or None
    except Exception as e:
        print("reverse_geocode:", e)
        return None


@app.post("/time-in")
async def time_in(
    authorization: str | None = Header(default=None),
    photo: UploadFile | None = File(None),
    latitude: float | None = Form(None),
    longitude: float | None = Form(None),
    readable_location_name: str | None = Form(None),
    capture_timestamp: str | None = Form(None),
):
    """Time In requires multipart: photo + GPS + capture timestamp (no manual time/location)."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        student_id, bad = resolve_student_id_number_or_403({}, authorization)
        if bad is not None:
            return bad

        if photo is None or not photo.filename:
            return JSONResponse(
                {
                    "error": "Photo is required. Tap Take Photo and allow camera access before Time In.",
                },
                status_code=400,
            )
        if latitude is None or longitude is None:
            return JSONResponse(
                {"error": "GPS location is required. Allow location access when taking your photo."},
                status_code=400,
            )
        try:
            lat = float(latitude)
            lon = float(longitude)
        except (TypeError, ValueError):
            return JSONResponse({"error": "Invalid GPS coordinates."}, status_code=400)
        if not (-90 <= lat <= 90 and -180 <= lon <= 180):
            return JSONResponse({"error": "GPS coordinates are out of range."}, status_code=400)

        try:
            capture_dt = _parse_capture_timestamp_iso(capture_timestamp or "")
        except ValueError as e:
            return JSONResponse({"error": str(e)}, status_code=400)

        now_dt = datetime.now(timezone.utc)
        skew_sec = abs((now_dt - capture_dt).total_seconds())
        if skew_sec > IMMERSION_CAPTURE_MAX_SKEW_MINUTES * 60:
            return JSONResponse(
                {
                    "error": f"Capture time must be within {IMMERSION_CAPTURE_MAX_SKEW_MINUTES} minutes of now. Take a new photo.",
                },
                status_code=400,
            )

        location_label = (readable_location_name or "").strip()
        if not location_label:
            location_label = _reverse_geocode_location(lat, lon) or ""
        if not location_label:
            return JSONResponse(
                {
                    "error": "Could not determine your location name. Allow GPS and try again outdoors or near a window.",
                },
                status_code=400,
            )

        existing = db_supabase.get_active_attendance(student_id)
        if existing:
            return JSONResponse(
                {"error": "Student already has an active session. Time Out first.", "attendance_id": existing.get("id")},
                status_code=409,
            )

        raw = await photo.read()
        content_type = (photo.content_type or "").lower()
        if content_type and not content_type.startswith("image/"):
            return JSONResponse({"error": "Only image files are allowed for Time In."}, status_code=400)
        if len(raw) > immersion_upload.MAX_PHOTO_BYTES:
            return JSONResponse({"error": "Photo is too large (max 6 MB)."}, status_code=400)

        photo_b64 = base64.standard_b64encode(raw).decode("ascii")

        now_iso = now_dt.isoformat()
        capture_iso = capture_dt.isoformat()
        row = db_supabase.insert_time_in_with_capture(
            student_id,
            now_iso,
            captured_photo_base64=photo_b64,
            latitude=lat,
            longitude=lon,
            readable_location_name=location_label,
            capture_timestamp=capture_iso,
        )
        return row
    except Exception as e:
        err_text = str(e)
        if "PGRST204" in err_text or "capture_timestamp" in err_text:
            return JSONResponse(
                {
                    "error": (
                        "Database setup incomplete: run backend/migrations/immersion_attendance_capture.sql "
                        "(and backend/migrations/immersion_attendance_photos_base64.sql if the error mentions "
                        "captured_photo_base64) in Supabase SQL Editor, wait a few seconds, then try Time In again."
                    ),
                },
                status_code=503,
            )
        return JSONResponse({"error": err_text}, status_code=502)


@app.post("/time-out")
async def time_out(
    authorization: str | None = Header(default=None),
    photo: UploadFile | None = File(None),
    latitude: float | None = Form(None),
    longitude: float | None = Form(None),
    readable_location_name: str | None = Form(None),
    capture_timestamp: str | None = Form(None),
):
    """Time Out requires multipart: photo + GPS + capture timestamp."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        student_id, bad = resolve_student_id_number_or_403({}, authorization)
        if bad is not None:
            return bad

        if photo is None or not photo.filename:
            return JSONResponse(
                {"error": "Photo is required. Tap Take Photo and allow camera access before Time Out."},
                status_code=400,
            )
        if latitude is None or longitude is None:
            return JSONResponse(
                {"error": "GPS location is required. Allow location access when taking your photo."},
                status_code=400,
            )
        try:
            lat = float(latitude)
            lon = float(longitude)
        except (TypeError, ValueError):
            return JSONResponse({"error": "Invalid GPS coordinates."}, status_code=400)
        if not (-90 <= lat <= 90 and -180 <= lon <= 180):
            return JSONResponse({"error": "GPS coordinates are out of range."}, status_code=400)

        try:
            capture_dt = _parse_capture_timestamp_iso(capture_timestamp or "")
        except ValueError as e:
            return JSONResponse({"error": str(e)}, status_code=400)

        now_dt = datetime.now(timezone.utc)
        if abs((now_dt - capture_dt).total_seconds()) > IMMERSION_CAPTURE_MAX_SKEW_MINUTES * 60:
            return JSONResponse(
                {
                    "error": f"Capture time must be within {IMMERSION_CAPTURE_MAX_SKEW_MINUTES} minutes of now. Take a new photo.",
                },
                status_code=400,
            )

        location_label = (readable_location_name or "").strip()
        if not location_label:
            location_label = _reverse_geocode_location(lat, lon) or ""
        if not location_label:
            return JSONResponse(
                {"error": "Could not determine your location name. Allow GPS and try again."},
                status_code=400,
            )

        active = db_supabase.get_active_attendance(student_id)
        if not active:
            return JSONResponse({"error": "No active Time In found for this student."}, status_code=400)

        raw = await photo.read()
        content_type = (photo.content_type or "").lower()
        if content_type and not content_type.startswith("image/"):
            return JSONResponse({"error": "Only image files are allowed for Time Out."}, status_code=400)
        if len(raw) > immersion_upload.MAX_PHOTO_BYTES:
            return JSONResponse({"error": "Photo is too large (max 6 MB)."}, status_code=400)

        photo_b64 = base64.standard_b64encode(raw).decode("ascii")

        now_iso = now_dt.isoformat()
        capture_iso = capture_dt.isoformat()
        updated = db_supabase.complete_time_out_with_capture(
            str(active["id"]),
            now_iso,
            time_out_photo_base64=photo_b64,
            latitude=lat,
            longitude=lon,
            readable_location_name=location_label,
            capture_timestamp=capture_iso,
        )
        return updated
    except Exception as e:
        err_text = str(e)
        if "PGRST204" in err_text and "time_out_" in err_text:
            return JSONResponse(
                {
                    "error": (
                        "Database setup incomplete: run backend/migrations/immersion_time_out_capture.sql "
                        "(and backend/migrations/immersion_attendance_photos_base64.sql if the error mentions "
                        "time_out_photo_base64) in Supabase SQL Editor, then try Time Out again."
                    ),
                },
                status_code=503,
            )
        return JSONResponse({"error": err_text}, status_code=502)


IMMERSION_QR_TTL_SECONDS = 20  # a bit longer than the 15s client refresh so an
# in-flight scan still lands even if the code rotated a moment ago


def _qr_signing_key() -> bytes:
    key = (
        os.getenv("SUPABASE_SERVICE_ROLE_KEY")
        or os.getenv("SUPABASE_KEY")
        or API_KEY
        or "learniq-fallback-qr-signing-key"
    ).strip()
    return key.encode("utf-8")


def _sign_immersion_qr_token(teacher_id_number: str, expires_at: int) -> str:
    msg = f"{teacher_id_number}:{expires_at}".encode("utf-8")
    sig = hmac.new(_qr_signing_key(), msg, hashlib.sha256).hexdigest()
    return f"IMMERSION:{teacher_id_number}:{expires_at}:{sig}"


def _verify_immersion_qr_token(scanned_code: str) -> tuple[str | None, str | None]:
    """Returns (teacher_id_number, error_message) — error_message is None on
    success. Rejects malformed codes, bad signatures, and expired codes (a
    saved screenshot stops working within IMMERSION_QR_TTL_SECONDS)."""
    parts = scanned_code.split(":")
    if len(parts) != 4 or parts[0].strip().upper() != "IMMERSION":
        return None, "That QR code is not a work immersion check-in code."
    _, teacher_id_number, expires_raw, sig = parts
    teacher_id_number = teacher_id_number.strip()
    if not teacher_id_number:
        return None, "That QR code is not a work immersion check-in code."
    try:
        expires_at = int(expires_raw)
    except ValueError:
        return None, "That QR code is not a work immersion check-in code."
    expected_sig = hmac.new(
        _qr_signing_key(), f"{teacher_id_number}:{expires_at}".encode("utf-8"), hashlib.sha256
    ).hexdigest()
    if not hmac.compare_digest(expected_sig, sig.strip()):
        return None, "That QR code is not a work immersion check-in code."
    if int(time.time()) > expires_at:
        return None, "This QR code has expired. Ask your workplace to refresh it, then scan again."
    return teacher_id_number, None


@app.get("/teacher/immersion/qr-token")
def teacher_immersion_qr_token_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
):
    """Short-lived signed token for the workplace check-in QR. The frontend
    re-fetches this every ~15s and re-renders the QR — anti-fraud: a
    photographed/saved code stops scanning within seconds, proving the
    student scanned it live at the workplace."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    expires_at = int(time.time()) + IMMERSION_QR_TTL_SECONDS
    return {
        "token": _sign_immersion_qr_token(tid, expires_at),
        "expires_at": expires_at,
        "ttl_seconds": IMMERSION_QR_TTL_SECONDS,
    }


@app.get("/teacher/immersion/recent-checkins")
def teacher_immersion_recent_checkins_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
    since: str | None = Query(default=None),
):
    """Recent Time In / Time Out events for this teacher's own Grade 12
    immersion students. The workplace QR display polls this so it can
    announce each scan out loud as it happens."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    try:
        events = db_supabase.list_recent_immersion_checkins_for_teacher(tid, since=since)
        return {"events": events, "server_time": datetime.now(timezone.utc).isoformat()}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _immersion_display_token(teacher_id_number: str) -> str:
    """Stable (non-expiring) signed token identifying a teacher, for the
    no-login workplace display page. A workplace host bookmarks a URL with
    this token instead of logging in — nothing sensitive is exposed beyond
    what that teacher's own workplace QR page already shows."""
    return hmac.new(
        _qr_signing_key(), f"display:{teacher_id_number}".encode("utf-8"), hashlib.sha256
    ).hexdigest()


def _resolve_teacher_from_display_token(teacher_id_number: str, token: str) -> bool:
    tid = str(teacher_id_number or "").strip()
    tok = str(token or "").strip()
    if not tid or not tok:
        return False
    expected = _immersion_display_token(tid)
    return hmac.compare_digest(expected, tok)


@app.get("/teacher/immersion/display-link")
def teacher_immersion_display_link_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
):
    """The teacher's own shareable, no-login link for the workplace display
    page — give this to the host company/supervisor once. It only ever
    shows the rotating check-in QR and this teacher's own recent scans."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    return {"teacher_id_number": tid, "token": _immersion_display_token(tid)}


@app.get("/public/immersion/qr-token")
def public_immersion_qr_token_endpoint(
    tid: str = Query(...),
    token: str = Query(...),
):
    """No-login variant of /teacher/immersion/qr-token — for the workplace
    display page, authenticated by the shareable display token instead of a
    teacher's own login session."""
    err = require_supabase()
    if err is not None:
        return err
    if not _resolve_teacher_from_display_token(tid, token):
        return JSONResponse({"error": "Invalid or expired link. Ask your teacher for a new one."}, status_code=403)
    expires_at = int(time.time()) + IMMERSION_QR_TTL_SECONDS
    return {
        "token": _sign_immersion_qr_token(tid, expires_at),
        "expires_at": expires_at,
        "ttl_seconds": IMMERSION_QR_TTL_SECONDS,
    }


@app.get("/public/immersion/recent-checkins")
def public_immersion_recent_checkins_endpoint(
    tid: str = Query(...),
    token: str = Query(...),
    since: str | None = Query(default=None),
):
    """No-login variant of /teacher/immersion/recent-checkins, for the
    workplace display page's live voice announcements."""
    err = require_supabase()
    if err is not None:
        return err
    if not _resolve_teacher_from_display_token(tid, token):
        return JSONResponse({"error": "Invalid or expired link. Ask your teacher for a new one."}, status_code=403)
    try:
        events = db_supabase.list_recent_immersion_checkins_for_teacher(tid, since=since)
        return {"events": events, "server_time": datetime.now(timezone.utc).isoformat()}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/student/immersion/qr-checkin")
async def student_immersion_qr_checkin_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    """Student scans the workplace QR code to Time In or Time Out —
    whichever applies given their current state. No photo/GPS: presence at
    the workplace is proven by scanning a code that only stays valid for
    ~20 seconds and is physically posted/shown at the site."""
    err = require_supabase()
    if err is not None:
        return err
    student_id, bad = resolve_student_id_number_or_403({}, authorization)
    if bad is not None:
        return bad
    scanned_code = str(body.get("scanned_code") or "").strip()
    teacher_id_number, token_error = _verify_immersion_qr_token(scanned_code)
    if token_error:
        return JSONResponse({"error": token_error}, status_code=400)
    if not db_supabase.teacher_can_view_student_immersion(teacher_id_number, student_id):
        return JSONResponse(
            {"error": "This workplace code is not linked to your immersion class."},
            status_code=403,
        )

    now_iso = datetime.now(timezone.utc).isoformat()
    try:
        active = db_supabase.get_active_attendance(student_id)
        if active:
            record = db_supabase.complete_time_out(str(active["id"]), now_iso)
            return {"action": "time_out", "record": record}
        record = db_supabase.insert_time_in(student_id, now_iso)
        return {"action": "time_in", "record": record}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/attendance-history")
def attendance_history(
    limit: int = Query(60, ge=1, le=200),
    authorization: str | None = Header(default=None),
):
    """Immersion dashboard: active clock-in + recent rows + sum of completed hours."""
    err = require_supabase()
    if err is not None:
        return err
    try:
        student_id, bad = resolve_student_id_number_or_403({}, authorization)
        if bad is not None:
            return bad
        rows = [db_supabase.enrich_attendance_row(r) for r in db_supabase.list_attendance_by_student(student_id)[:limit]]
        active = db_supabase.enrich_attendance_row(db_supabase.get_active_attendance(student_id))
        total = 0.0
        for r in rows:
            th = r.get("total_hours")
            if th is not None:
                try:
                    total += float(th)
                except (TypeError, ValueError):
                    pass
        return {
            "active": active,
            "history": rows,
            "total_hours_rendered": round(total, 2),
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/submit-journal")
async def submit_journal(
    body: dict | None = Body(default=None),
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    try:
        payload = body if isinstance(body, dict) else {}
        student_id, bad = resolve_student_id_number_or_403(payload, authorization)
        if bad is not None:
            return bad
        journal_body = (payload.get("body") or "").strip()
        if not journal_body:
            return JSONResponse({"error": "body is required."}, status_code=400)

        attendance_id = payload.get("attendance_id") or payload.get("immersion_log_id")
        if attendance_id is not None:
            attendance_id = str(attendance_id)
        else:
            records = db_supabase.list_attendance_by_student(student_id)
            if not records:
                return JSONResponse({"error": "No attendance record found for this student."}, status_code=400)
            attendance_id = str(records[0]["id"])

        row = db_supabase.insert_journal_linked(student_id, attendance_id, journal_body)
        return row
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/attendance/{student_id}")
def get_attendance(student_id: str, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    allowed, _, bad = _can_view_student_data(authorization, student_id)
    if not allowed:
        return bad
    try:
        return db_supabase.list_attendance_by_student(student_id)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/journals/{student_id}")
def get_journals(student_id: str, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    if caller_idn != student_id:
        caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
        if str((caller_prof or {}).get("role") or "").strip().lower() != "admin":
            return JSONResponse({"error": "You don't have access to this student's journals."}, status_code=403)
    try:
        return db_supabase.list_journals_for_student(student_id)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/attendance")
async def attendance(body: dict, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    try:
        sid, bad = resolve_student_id_number_or_403(body, authorization)
        if bad is not None:
            return bad
        ev = (body.get("event_type") or "").strip().lower()
        if not sid or ev not in ("time_in", "time_out"):
            return JSONResponse(
                {"error": "student_id_number and event_type (time_in | time_out) are required."},
                status_code=400,
            )
        return db_supabase.insert_attendance(sid, ev)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/attendance")
def attendance_list(student_id_number: str, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    allowed, _, bad = _can_view_student_data(authorization, student_id_number)
    if not allowed:
        return bad
    try:
        return db_supabase.list_attendance_for_student(student_id_number)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/journals")
async def journals_create(body: dict, authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    try:
        sid, bad = resolve_student_id_number_or_403(body, authorization)
        if bad is not None:
            return bad
        text_body = (body.get("body") or body.get("journal_text") or "").strip()
        if not text_body:
            return JSONResponse({"error": "body (or journal_text) is required."}, status_code=400)
        return db_supabase.insert_journal(sid, text_body, entry_date=body.get("entry_date"))
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/journals")
def journals_list(
    student_id_number: str | None = None,
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    try:
        sid, bad = resolve_student_id_number_or_403(
            {"student_id_number": student_id_number} if student_id_number else {},
            authorization,
        )
        if bad is not None:
            return bad
        return db_supabase.list_journals_for_student(sid)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


# ============================================================
# Student Gradecard endpoints
# ============================================================

@app.get("/grading-periods")
def list_grading_periods_endpoint():
    err = require_supabase()
    if err is not None:
        return err
    try:
        rows = db_supabase.list_grading_periods()
        current = db_supabase.get_current_grading_period() or {}
        return {"periods": rows, "current_id": current.get("id")}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/grading-periods/current")
def current_grading_period_endpoint():
    err = require_supabase()
    if err is not None:
        return err
    try:
        period = db_supabase.get_current_grading_period()
        if not period:
            return JSONResponse(
                {"error": "No grading period configured. Run the gradecard migration SQL."},
                status_code=404,
            )
        return period
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/gradecard/strands")
def teacher_gradecard_strands_endpoint(
    teacher_id_number: str = Query(...),
    authorization: str | None = Header(default=None),
):
    """Strand tiles with student counts (teacher's enrolled students only)."""
    err = require_supabase()
    if err is not None:
        return err
    tid = str(teacher_id_number or "").strip()
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    allowed, bad = _can_view_teacher_data(authorization, tid)
    if not allowed:
        return bad
    try:
        strands = db_supabase.list_gradecard_strands_for_teacher(tid)
        return {"strands": strands, "count": len(strands)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _resolve_teacher_or_admin_id(
    authorization: str | None,
    teacher_id_number: str | None,
) -> tuple[str | None, JSONResponse | None]:
    """Bearer token must be teacher (own id) or admin (optional teacher_id_number override)."""
    token_idn = student_id_number_from_authorization(authorization)
    if not token_idn:
        return None, JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        prof = db_supabase.get_profile_by_id_number(token_idn)
    except Exception:
        prof = None
    role = str((prof or {}).get("role") or "").strip().lower()
    if role == "admin":
        tid = str(teacher_id_number or token_idn).strip()
        return (tid or None), (
            JSONResponse({"error": "teacher_id_number is required for admin."}, status_code=400)
            if not tid
            else None
        )
    if role == "teacher":
        q_tid = str(teacher_id_number or "").strip()
        if q_tid and q_tid != token_idn:
            return None, JSONResponse(
                {"error": "teacher_id_number does not match signed-in teacher."},
                status_code=403,
            )
        return token_idn, None
    return None, JSONResponse({"error": "Teacher or admin access only."}, status_code=403)


def _resolve_admin_id(authorization: str | None) -> tuple[str | None, JSONResponse | None]:
    """Bearer token must belong to an admin. Returns (admin's own id_number, error)."""
    token_idn = student_id_number_from_authorization(authorization)
    if not token_idn:
        return None, JSONResponse({"error": "Sign in required."}, status_code=401)
    try:
        prof = db_supabase.get_profile_by_id_number(token_idn)
    except Exception:
        prof = None
    role = str((prof or {}).get("role") or "").strip().lower()
    if role != "admin":
        return None, JSONResponse({"error": "Admin access only."}, status_code=403)
    return token_idn, None


@app.get("/sections")
def list_sections_endpoint(
    grade_level: str | None = Query(default=None),
    strand: str | None = Query(default=None),
):
    """Fixed section list — readable by any registration form (admin panel
    today; not sensitive data, just section names)."""
    err = require_supabase()
    if err is not None:
        return err
    sections = db_supabase.list_sections(grade_level=grade_level, strand=strand)
    return {"sections": sections}


@app.post("/admin/sections")
async def admin_create_section_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        section = db_supabase.create_section(
            body.get("name"), body.get("grade_level"), body.get("strand")
        )
        return {"section": section}
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.delete("/admin/sections/{section_id}")
async def admin_delete_section_endpoint(
    section_id: str,
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        ok = db_supabase.delete_section(section_id)
        if not ok:
            return JSONResponse({"error": "Section not found."}, status_code=404)
        return {"ok": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/admin/students-without-section")
def admin_students_without_section_endpoint(
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        students = db_supabase.list_students_without_section()
        return {"students": students, "count": len(students)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/admin/students/section")
async def admin_set_student_section_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    err = require_supabase()
    if err is not None:
        return err
    _, bad = _resolve_admin_id(authorization)
    if bad is not None:
        return bad
    try:
        profile = db_supabase.admin_set_student_section(
            body.get("student_id_number"), body.get("section")
        )
        if not profile:
            return JSONResponse({"error": "Student not found."}, status_code=404)
        return {"profile": db_supabase.serialize_public_profile(profile)}
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/immersion/strands")
def teacher_immersion_strands_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
):
    """Grade 12 immersion: strand tiles with student counts (teacher roster)."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    try:
        strands = db_supabase.list_gradecard_strands_for_teacher(tid)
        return {"strands": strands, "grade_level": "12", "count": len(strands)}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/immersion/students")
def teacher_immersion_students_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
    strand: str = Query(...),
    q: str | None = Query(default=None),
):
    """Grade 12 students in a strand enrolled in the teacher's subjects."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    st = str(strand or "").strip()
    if not st:
        return JSONResponse({"error": "strand is required"}, status_code=400)
    try:
        students = db_supabase.list_gradecard_students_for_teacher(
            tid, st, search=q, grade_level="12"
        )
        students = db_supabase.enrich_students_with_immersion_status(students)
        return {
            "strand": st,
            "grade_level": "12",
            "students": students,
            "count": len(students),
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


def _bearer_token(authorization: str | None, access_token: str | None = None) -> str | None:
    if authorization and authorization.lower().startswith("bearer "):
        t = authorization.split(" ", 1)[1].strip()
        if t:
            return t
    return (access_token or "").strip() or None


@app.get("/teacher/immersion/attendance-photo")
def teacher_immersion_attendance_photo_endpoint(
    student_id_number: str = Query(...),
    attendance_id: str = Query(...),
    kind: str = Query(...),
    authorization: str | None = Header(default=None),
    access_token: str | None = Query(default=None),
    teacher_id_number: str | None = Query(default=None),
):
    """Stream verification photo bytes (for img src; pass access_token query if needed)."""
    err = require_supabase()
    if err is not None:
        return err
    token = _bearer_token(authorization, access_token)
    auth_header = f"Bearer {token}" if token else None
    tid, bad = _resolve_teacher_or_admin_id(auth_header, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    sid = str(student_id_number or "").strip()
    aid = str(attendance_id or "").strip()
    photo_kind = str(kind or "").strip().lower()
    if photo_kind not in ("time_in", "time_out"):
        return JSONResponse({"error": "kind must be time_in or time_out"}, status_code=400)
    if not sid or not aid:
        return JSONResponse(
            {"error": "student_id_number and attendance_id are required"},
            status_code=400,
        )
    try:
        if not db_supabase.teacher_can_view_student_immersion(tid, sid):
            return JSONResponse({"error": "Access denied."}, status_code=403)
        row = db_supabase.get_attendance_log_by_id(aid)
        if not row or not db_supabase._attendance_belongs_to_student(row, sid):
            return JSONResponse({"error": "Attendance session not found."}, status_code=404)
        loaded = db_supabase.read_attendance_photo_bytes(row, photo_kind)
        if not loaded:
            return JSONResponse({"error": "Photo not found for this session."}, status_code=404)
        data, mime = loaded
        return Response(
            content=data,
            media_type=mime,
            headers={"Cache-Control": "private, max-age=300"},
        )
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/immersion/attendance-session")
def teacher_immersion_attendance_session_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
    student_id_number: str = Query(...),
    attendance_id: str = Query(...),
):
    """One clock session with Time In / Time Out verification photos (teacher review)."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    sid = str(student_id_number or "").strip()
    aid = str(attendance_id or "").strip()
    if not sid or not aid:
        return JSONResponse(
            {"error": "student_id_number and attendance_id are required"},
            status_code=400,
        )
    try:
        session = db_supabase.get_teacher_immersion_attendance_session(tid, sid, aid)
        return {"session": session}
    except PermissionError as e:
        return JSONResponse({"error": str(e)}, status_code=403)
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/immersion/student-overview")
def teacher_immersion_student_overview_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
    student_id_number: str = Query(...),
    limit: int = Query(default=120, ge=1, le=200),
):
    """Immersion attendance, hours, and journals for one Grade 12 student."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    try:
        return db_supabase.build_teacher_immersion_student_overview(
            tid, sid, limit=limit
        )
    except PermissionError as e:
        return JSONResponse({"error": str(e)}, status_code=403)
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/teacher/class-attendance/start")
async def teacher_class_attendance_start_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    """Open today's class attendance for a subject (students can photo check-in)."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(
        authorization, body.get("teacher_id_number")
    )
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    subject_id = str(body.get("subject_id") or "").strip()
    if not subject_id:
        return JSONResponse({"error": "subject_id is required"}, status_code=400)
    teacher_lat = body.get("teacher_latitude")
    teacher_lon = body.get("teacher_longitude")
    teacher_loc = str(body.get("teacher_start_location_name") or "").strip() or None
    try:
        lat = float(teacher_lat) if teacher_lat is not None else None
        lon = float(teacher_lon) if teacher_lon is not None else None
    except (TypeError, ValueError):
        lat, lon = None, None
    try:
        session = db_supabase.start_class_attendance_session(
            subject_id,
            tid,
            teacher_latitude=lat,
            teacher_longitude=lon,
            teacher_start_location_name=teacher_loc,
        )
        geofence = db_supabase._geofence_for_api(subject_id, session)
        return {"session": session, "geofence": geofence}
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/teacher/class-attendance/end")
async def teacher_class_attendance_end_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    """Close class attendance session for today."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(
        authorization, body.get("teacher_id_number")
    )
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    session_id = str(body.get("session_id") or "").strip()
    if not session_id:
        return JSONResponse({"error": "session_id is required"}, status_code=400)
    try:
        session = db_supabase.end_class_attendance_session(session_id, tid)
        return {"session": session}
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/class-attendance/live")
def teacher_class_attendance_live_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
    subject_id: str = Query(...),
):
    """Live roster for polling while session is open (or today's final roster)."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    sid = str(subject_id or "").strip()
    if not sid:
        return JSONResponse({"error": "subject_id is required"}, status_code=400)
    try:
        return db_supabase.build_class_attendance_live(sid, tid)
    except PermissionError as e:
        return JSONResponse({"error": str(e)}, status_code=403)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/student/class-attendance/status")
def student_class_attendance_status_endpoint(
    authorization: str | None = Header(default=None),
    subject_id: str = Query(...),
):
    """Whether class attendance is open for this subject and if the student checked in."""
    err = require_supabase()
    if err is not None:
        return err
    student_id, bad = resolve_student_id_number_or_403({}, authorization)
    if bad is not None:
        return bad
    sid = str(subject_id or "").strip()
    if not sid:
        return JSONResponse({"error": "subject_id is required"}, status_code=400)
    prof = db_supabase.get_profile_by_id_number(student_id)
    if not prof:
        return JSONResponse({"error": "Student profile not found"}, status_code=404)
    student_uuid = str(prof.get("id") or "")
    try:
        return db_supabase.get_student_class_attendance_status(
            student_uuid, student_id, sid
        )
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/student/class-attendance/check-in")
async def student_class_attendance_check_in_endpoint(
    authorization: str | None = Header(default=None),
    subject_id: str = Form(...),
    photo: UploadFile | None = File(None),
    latitude: float | None = Form(None),
    longitude: float | None = Form(None),
    readable_location_name: str | None = Form(None),
    capture_timestamp: str | None = Form(None),
):
    """Photo check-in for class attendance (same capture rules as immersion Time In)."""
    err = require_supabase()
    if err is not None:
        return err
    student_id, bad = resolve_student_id_number_or_403({}, authorization)
    if bad is not None:
        return bad
    sid = str(subject_id or "").strip()
    if not sid:
        return JSONResponse({"error": "subject_id is required"}, status_code=400)
    prof = db_supabase.get_profile_by_id_number(student_id)
    if not prof:
        return JSONResponse({"error": "Student profile not found"}, status_code=404)
    student_uuid = str(prof.get("id") or "")

    if photo is None or not photo.filename:
        return JSONResponse(
            {"error": "Photo is required. Take a photo before checking in."},
            status_code=400,
        )
    if latitude is None or longitude is None:
        return JSONResponse(
            {"error": "GPS location is required. Allow location when taking your photo."},
            status_code=400,
        )
    try:
        lat = float(latitude)
        lon = float(longitude)
    except (TypeError, ValueError):
        return JSONResponse({"error": "Invalid GPS coordinates."}, status_code=400)

    try:
        capture_dt = _parse_capture_timestamp_iso(capture_timestamp or "")
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)

    now_dt = datetime.now(timezone.utc)
    if abs((now_dt - capture_dt).total_seconds()) > IMMERSION_CAPTURE_MAX_SKEW_MINUTES * 60:
        return JSONResponse(
            {
                "error": f"Capture time must be within {IMMERSION_CAPTURE_MAX_SKEW_MINUTES} minutes of now.",
            },
            status_code=400,
        )

    location_label = (readable_location_name or "").strip()
    if not location_label:
        location_label = _reverse_geocode_location(lat, lon) or ""
    if not location_label:
        return JSONResponse(
            {"error": "Could not determine your location. Allow GPS and try again."},
            status_code=400,
        )

    raw = await photo.read()
    content_type = (photo.content_type or "").lower()
    if content_type and not content_type.startswith("image/"):
        return JSONResponse({"error": "Only image files are allowed."}, status_code=400)
    if len(raw) > immersion_upload.MAX_PHOTO_BYTES:
        return JSONResponse({"error": "Photo is too large (max 6 MB)."}, status_code=400)

    photo_b64 = base64.standard_b64encode(raw).decode("ascii")
    try:
        record = db_supabase.insert_class_attendance_checkin(
            student_uuid,
            student_id,
            sid,
            time_in_iso=now_dt.isoformat(),
            captured_photo_base64=photo_b64,
            latitude=lat,
            longitude=lon,
            readable_location_name=location_label,
            capture_timestamp=capture_dt.isoformat(),
        )
        return {"record": record}
    except PermissionError as e:
        return JSONResponse({"error": str(e)}, status_code=403)
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/teacher/class-attendance/scan")
async def teacher_class_attendance_scan_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    """Teacher scans a student's QR code (their id_number) with the class
    camera to mark them present — no photo/GPS step for the student."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, body.get("teacher_id_number"))
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    subject_id = str(body.get("subject_id") or "").strip()
    if not subject_id:
        return JSONResponse({"error": "subject_id is required"}, status_code=400)
    scanned_id_number = str(body.get("scanned_id_number") or "").strip()
    if not scanned_id_number:
        return JSONResponse({"error": "scanned_id_number is required"}, status_code=400)

    prof = db_supabase.get_profile_by_id_number(scanned_id_number)
    if not prof or str(prof.get("role") or "").lower() != "student":
        return JSONResponse({"error": "QR code does not match a student account."}, status_code=404)
    student_uuid = str(prof.get("id") or "")
    try:
        record = db_supabase.insert_class_attendance_checkin_via_qr(
            student_uuid,
            scanned_id_number,
            subject_id,
            time_in_iso=datetime.now(timezone.utc).isoformat(),
        )
        return {
            "record": record,
            "student_id_number": scanned_id_number,
            "student_name": db_supabase.profile_display_name(prof) or scanned_id_number,
        }
    except PermissionError as e:
        return JSONResponse({"error": str(e)}, status_code=403)
    except ValueError as e:
        msg = str(e)
        status_code = 409 if "already checked in" in msg.lower() else 400
        return JSONResponse({"error": msg}, status_code=status_code)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


CLASS_ATTENDANCE_QR_TTL_SECONDS = 15  # matches the on-screen rotation cadence


def _sign_class_attendance_qr_token(subject_id: str, expires_at: int) -> str:
    msg = f"{subject_id}:{expires_at}".encode("utf-8")
    sig = hmac.new(_qr_signing_key(), msg, hashlib.sha256).hexdigest()
    return f"CLASSATT:{subject_id}:{expires_at}:{sig}"


def _verify_class_attendance_qr_token(scanned_code: str) -> tuple[str | None, str | None]:
    """Returns (subject_id, error_message) — error_message is None on
    success. Same anti-fraud model as the Immersion workplace QR: a saved
    screenshot stops scanning within CLASS_ATTENDANCE_QR_TTL_SECONDS."""
    parts = scanned_code.split(":")
    if len(parts) != 4 or parts[0].strip().upper() != "CLASSATT":
        return None, "That QR code is not a class attendance check-in code."
    _, subject_id, expires_raw, sig = parts
    subject_id = subject_id.strip()
    if not subject_id:
        return None, "That QR code is not a class attendance check-in code."
    try:
        expires_at = int(expires_raw)
    except ValueError:
        return None, "That QR code is not a class attendance check-in code."
    expected_sig = hmac.new(
        _qr_signing_key(), f"{subject_id}:{expires_at}".encode("utf-8"), hashlib.sha256
    ).hexdigest()
    if not hmac.compare_digest(expected_sig, sig.strip()):
        return None, "That QR code is not a class attendance check-in code."
    if int(time.time()) > expires_at:
        return None, "This QR code has expired. Ask your teacher to refresh it, then scan again."
    return subject_id, None


@app.get("/teacher/class-attendance/qr-token")
def teacher_class_attendance_qr_token_endpoint(
    authorization: str | None = Header(default=None),
    teacher_id_number: str | None = Query(default=None),
    subject_id: str = Query(...),
):
    """Short-lived signed token for the class attendance QR the teacher
    displays. The frontend re-fetches this every ~15s and re-renders the QR
    — a photographed/saved code stops scanning within seconds, proving the
    student scanned it live in class."""
    err = require_supabase()
    if err is not None:
        return err
    tid, bad = _resolve_teacher_or_admin_id(authorization, teacher_id_number)
    if bad is not None:
        return bad
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    sid = str(subject_id or "").strip()
    if not sid:
        return JSONResponse({"error": "subject_id is required"}, status_code=400)
    expires_at = int(time.time()) + CLASS_ATTENDANCE_QR_TTL_SECONDS
    return {
        "token": _sign_class_attendance_qr_token(sid, expires_at),
        "expires_at": expires_at,
        "ttl_seconds": CLASS_ATTENDANCE_QR_TTL_SECONDS,
    }


@app.post("/student/class-attendance/qr-checkin")
async def student_class_attendance_qr_checkin_endpoint(
    authorization: str | None = Header(default=None),
    body: dict = Body(...),
):
    """Student scans the teacher's rotating QR to check themselves in or
    out — first scan = time in, second scan (before the session ends) =
    time out, mirroring Immersion's time-in/time-out. No photo/GPS: the QR
    itself expires in ~15s, so only someone physically looking at the
    teacher's screen right now can scan it."""
    err = require_supabase()
    if err is not None:
        return err
    student_id, bad = resolve_student_id_number_or_403({}, authorization)
    if bad is not None:
        return bad
    scanned_code = str(body.get("scanned_code") or "").strip()
    subject_id, token_error = _verify_class_attendance_qr_token(scanned_code)
    if token_error:
        return JSONResponse({"error": token_error}, status_code=400)

    prof = db_supabase.get_profile_by_id_number(student_id)
    if not prof or str(prof.get("role") or "").lower() != "student":
        return JSONResponse({"error": "Student profile not found."}, status_code=404)
    student_uuid = str(prof.get("id") or "")
    try:
        action, record = db_supabase.class_attendance_qr_toggle(
            student_uuid,
            student_id,
            subject_id,
            now_iso=datetime.now(timezone.utc).isoformat(),
        )
        return {
            "action": action,
            "record": record,
            "student_id_number": student_id,
            "student_name": db_supabase.profile_display_name(prof) or student_id,
        }
    except PermissionError as e:
        return JSONResponse({"error": str(e)}, status_code=403)
    except ValueError as e:
        msg = str(e)
        status_code = 409 if "already checked out" in msg.lower() else 400
        return JSONResponse({"error": msg}, status_code=status_code)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/teacher/gradecard/students")
def teacher_gradecard_students_endpoint(
    teacher_id_number: str = Query(...),
    strand: str = Query(...),
    q: str | None = Query(default=None),
    grade_level: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
):
    """Students in a strand enrolled in the teacher's subjects."""
    err = require_supabase()
    if err is not None:
        return err
    tid = str(teacher_id_number or "").strip()
    if not tid:
        return JSONResponse({"error": "teacher_id_number is required"}, status_code=400)
    allowed, bad = _can_view_teacher_data(authorization, tid)
    if not allowed:
        return bad
    st = str(strand or "").strip()
    if not st:
        return JSONResponse({"error": "strand is required"}, status_code=400)
    gl = str(grade_level or "").strip() or None
    if gl and gl not in ("11", "12"):
        return JSONResponse(
            {"error": "grade_level must be 11 or 12"}, status_code=400
        )
    try:
        students = db_supabase.list_gradecard_students_for_teacher(
            tid, st, search=q, grade_level=gl
        )
        return {
            "strand": st,
            "grade_level": gl,
            "students": students,
            "count": len(students),
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/gradecard")
def get_gradecard_endpoint(
    student_id_number: str = Query(...),
    period_id: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    allowed, _, bad = _can_view_student_data(authorization, sid)
    if not allowed:
        return bad
    try:
        data = db_supabase.build_full_gradecard(sid, period_id)
        return data
    except RuntimeError as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except Exception as e:
        import traceback
        print("GRADECARD ERROR:", repr(e))
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/student-grades")
async def save_student_grade_endpoint(body: dict = Body(...), authorization: str | None = Header(default=None)):
    """Teacher saves / updates a per-subject grade row."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    caller_role = str((caller_prof or {}).get("role") or "").strip().lower()
    if caller_role not in ("teacher", "admin"):
        return JSONResponse({"error": "Only teachers or admins can save grades."}, status_code=403)

    payload = body if isinstance(body, dict) else {}

    student_id_number = str(payload.get("student_id_number") or "").strip()
    subject_id = payload.get("subject_id")
    period_id = payload.get("grading_period_id")

    if not student_id_number or not subject_id:
        return JSONResponse(
            {"error": "student_id_number and subject_id are required"},
            status_code=400,
        )

    if not period_id:
        cur = db_supabase.get_current_grading_period() or {}
        period_id = cur.get("id")
    if not period_id:
        return JSONResponse({"error": "No grading_period_id and no current period configured."}, status_code=400)

    student_uuid = db_supabase.profile_uuid_for_id_number(student_id_number)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)

    profile = db_supabase.get_profile_by_id_number(student_id_number) or {}
    # A teacher can only sign a grade as themselves — the token identity, not
    # whatever teacher_id_number the request body claims.
    partial: dict = {
        "teacher_id_number": caller_idn if caller_role == "teacher" else payload.get("teacher_id_number"),
        "written_work_score": payload.get("written_work_score"),
        "performance_task_score": payload.get("performance_task_score"),
        "quarterly_assessment_score": payload.get("quarterly_assessment_score"),
        "attendance_percent": payload.get("attendance_percent"),
        "final_grade": payload.get("final_grade"),
        "final_is_manual": payload.get("final_is_manual"),
        "remarks": payload.get("remarks"),
        "teacher_comments": payload.get("teacher_comments"),
        "finalized_at": payload.get("finalized_at"),
    }
    if payload.get("final_is_override") is True and "final_is_manual" not in payload:
        partial["final_is_manual"] = True

    try:
        db_payload = db_supabase.merge_and_compute_student_grade(
            student_uuid,
            profile.get("strand"),
            str(subject_id),
            period_id,
            partial,
        )
        row = db_supabase.upsert_student_grade(db_payload)
        return {"ok": True, "grade": row}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/gradecards")
async def save_gradecard_endpoint(body: dict = Body(...), authorization: str | None = Header(default=None)):
    """Adviser/Admin saves the top-level gradecard summary."""
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    if str((caller_prof or {}).get("role") or "").strip().lower() not in ("teacher", "admin"):
        return JSONResponse({"error": "Only teachers or admins can save gradecards."}, status_code=403)
    payload = body if isinstance(body, dict) else {}

    student_id_number = str(payload.get("student_id_number") or "").strip()
    period_id = payload.get("grading_period_id")

    if not student_id_number:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)

    period = db_supabase.get_grading_period(period_id)
    if not period:
        return JSONResponse({"error": "No grading period available"}, status_code=400)

    student = db_supabase.get_profile_by_id_number(student_id_number)
    if not student:
        return JSONResponse({"error": "Student not found"}, status_code=404)

    db_payload = {
        "student_id": str(student["id"]),
        "grading_period_id": period["id"],
        "general_average": payload.get("general_average"),
        "standing": payload.get("standing"),
        "conduct": payload.get("conduct"),
        "days_present": payload.get("days_present"),
        "days_absent": payload.get("days_absent"),
        "times_tardy": payload.get("times_tardy"),
        "adviser_comments": payload.get("adviser_comments"),
        "reference_no": payload.get("reference_no"),
        "finalized_at": payload.get("finalized_at"),
    }

    try:
        row = db_supabase.upsert_gradecard(db_payload)
        return {"ok": True, "gradecard": row}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.get("/enrollments")
def list_enrollments_endpoint(
    student_id_number: str = Query(...),
    period_id: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
):
    err = require_supabase()
    if err is not None:
        return err
    sid = str(student_id_number or "").strip()
    if not sid:
        return JSONResponse({"error": "student_id_number is required"}, status_code=400)
    allowed, _, bad = _can_view_student_data(authorization, sid)
    if not allowed:
        return bad
    student_uuid = db_supabase.profile_uuid_for_id_number(sid)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)
    if not period_id:
        cur = db_supabase.get_current_grading_period() or {}
        period_id = cur.get("id")
    try:
        rows = db_supabase.list_enrollments_for_student(student_uuid, period_id)
        return {"enrollments": rows}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/enrollments")
async def upsert_enrollment_endpoint(body: dict = Body(...), authorization: str | None = Header(default=None)):
    err = require_supabase()
    if err is not None:
        return err
    caller_idn = student_id_number_from_authorization(authorization)
    if not caller_idn:
        return JSONResponse({"error": "Sign in required."}, status_code=401)
    caller_prof = db_supabase.get_profile_by_id_number(caller_idn)
    if str((caller_prof or {}).get("role") or "").strip().lower() not in ("teacher", "admin"):
        return JSONResponse({"error": "Only teachers or admins can manage enrollments."}, status_code=403)
    payload = body if isinstance(body, dict) else {}

    student_id_number = str(payload.get("student_id_number") or "").strip()
    subject_id = payload.get("subject_id")
    if not student_id_number or not subject_id:
        return JSONResponse(
            {"error": "student_id_number and subject_id are required"},
            status_code=400,
        )

    student_uuid = db_supabase.profile_uuid_for_id_number(student_id_number)
    if not student_uuid:
        return JSONResponse({"error": "Student not found"}, status_code=404)

    period_id = payload.get("grading_period_id")
    if not period_id:
        cur = db_supabase.get_current_grading_period() or {}
        period_id = cur.get("id")

    try:
        row = db_supabase.upsert_enrollment(
            student_uuid=student_uuid,
            subject_id=str(subject_id),
            teacher_id_number=payload.get("teacher_id_number"),
            grading_period_id=period_id,
        )
        return {"ok": True, "enrollment": row}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


@app.post("/submit-activity")
async def submit_activity_endpoint(
    body: dict = Body(...),
    authorization: str | None = Header(default=None),
):
    """Student submits an essay response or flashcard review log."""
    err = require_supabase()
    if err is not None:
        return err
    payload = body if isinstance(body, dict) else {}

    student_id_number, bad = resolve_student_id_number_or_403(payload, authorization)
    if bad is not None:
        return bad

    lesson_id = payload.get("lesson_id") or payload.get("file_id")
    if not lesson_id:
        return JSONResponse({"error": "lesson_id is required"}, status_code=400)

    student_uuid = db_supabase.profile_uuid_for_id_number(student_id_number)
    if not student_uuid:
        return JSONResponse({"error": "Student profile not found"}, status_code=404)

    try:
        row = db_supabase.insert_activity_attempt({
            "student_id": student_uuid,
            "lesson_id": lesson_id,
            "activity_index": payload.get("activity_index") or 0,
            "activity_type": payload.get("activity_type") or "essay",
            "response": payload.get("response"),
            "score": payload.get("score"),
        })
        return {"ok": True, "attempt": row}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=502)


UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
LESSON_UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
(UPLOADS_DIR / "immersion").mkdir(parents=True, exist_ok=True)
app.mount("/uploads", StaticFiles(directory=str(UPLOADS_DIR)), name="uploads")

if FRONTEND_DIR.exists():
    app.mount("/", StaticFiles(directory=str(FRONTEND_DIR), html=True), name="frontend")
